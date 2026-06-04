"""
文档知识管理平台 - 全功能版本
Flask 后端主程序
"""
import os
import json
import uuid
import base64
import hashlib
import sqlite3
import datetime
import shutil
import re
import io
import traceback
from pathlib import Path
from functools import wraps
from threading import Lock
import numpy as np

from flask import (
    Flask, render_template, request, jsonify, send_file,
    session, g, url_for, redirect
)
from flask_cors import CORS
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from dotenv import load_dotenv

# ==================== Markdown 转 HTML（无需第三方库） ====================
def markdown_to_html(text):
    """简单的 Markdown 转 HTML，支持 AI 分析报告常用语法"""
    if not text:
        return ''
    
    # 保护数学公式（$$...$$ 块级、$...$ 行内），避免被 markdown 行解析破坏
    import re as _re
    math_blocks = []
    def _protect_math(m):
        math_blocks.append(m.group(0))
        return f'\x00MATHBLOCK{len(math_blocks)-1}\x00'
    text = _re.sub(r'\$\$[\s\S]+?\$\$', _protect_math, text)  # 多行块级公式
    text = _re.sub(r'\$[^$\n]+?\$', _protect_math, text)       # 行内公式
    
    lines = text.split('\n')
    html = []
    i = 0
    in_code_block = False
    code_buffer = []
    
    while i < len(lines):
        line = lines[i]
        
        # 代码块
        if line.strip().startswith('```'):
            if in_code_block:
                html.append(f'<pre><code>{"".join(code_buffer)}</code></pre>')
                code_buffer = []
                in_code_block = False
            else:
                in_code_block = True
            i += 1
            continue
        
        if in_code_block:
            code_buffer.append(line + '\n')
            i += 1
            continue
        
        stripped = line.strip()
        
        # 空行
        if not stripped:
            html.append('')
            i += 1
            continue
        
        # 标题
        if stripped.startswith('### '):
            html.append(f'<h3>{inline_md(stripped[4:])}</h3>')
        elif stripped.startswith('## '):
            html.append(f'<h2>{inline_md(stripped[3:])}</h2>')
        elif stripped.startswith('# '):
            html.append(f'<h1>{inline_md(stripped[2:])}</h1>')
        # 无序列表
        elif stripped.startswith('- ') or stripped.startswith('* '):
            html.append(f'<li>{inline_md(stripped[2:])}</li>')
        # 有序列表
        elif stripped[0].isdigit() and '. ' in stripped[:4]:
            html.append(f'<li>{inline_md(stripped.split(". ", 1)[1])}</li>')
        # 引用
        elif stripped.startswith('> '):
            html.append(f'<blockquote>{inline_md(stripped[2:])}</blockquote>')
        # 分隔线
        elif stripped in ('---', '***', '___'):
            html.append('<hr>')
        # 普通段落
        else:
            html.append(f'<p>{inline_md(stripped)}</p>')
        
        i += 1
    
    if in_code_block:
        html.append(f'<pre><code>{"".join(code_buffer)}</code></pre>')
    
    result = '\n'.join(html)
    
    # 将连续的 <li> 包裹成 <ul> 或 <ol>
    result = _wrap_lists(result)
    
    # 恢复数学公式（前端 KaTeX 负责渲染）
    for i, formula in enumerate(math_blocks):
        result = result.replace(f'\x00MATHBLOCK{i}\x00', formula)
    
    return result

def inline_md(text):
    """处理行内 Markdown 语法"""
    import re
    # 保护数学公式，避免 LaTeX 符号被 markdown 语法误处理（如 ** ^ _ 等）
    math_placeholders = []
    def _protect_math(m):
        math_placeholders.append(m.group(0))
        return f'\x00MATH{len(math_placeholders)-1}\x00'
    text = re.sub(r'\$\$[\s\S]+?\$\$', _protect_math, text)  # 块级公式 $$...$$
    text = re.sub(r'\$[^$\n]+?\$', _protect_math, text)       # 行内公式 $...$
    
    # 加粗 **text**
    text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)
    # 斜体 *text*
    text = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'<em>\1</em>', text)
    # 行内代码 `code`
    text = re.sub(r'`([^`]+)`', r'<code>\1</code>', text)
    # 链接 [text](url)
    text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'<a href="\2" target="_blank">\1</a>', text)
    # 图片 ![alt](url)
    text = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', r'<img src="\2" alt="\1" style="max-width:100%">', text)
    
    # 恢复数学公式（前端 KaTeX 负责渲染）
    for i, formula in enumerate(math_placeholders):
        text = text.replace(f'\x00MATH{i}\x00', formula)
    return text

def _wrap_lists(html):
    """将连续的 <li> 包裹成列表"""
    import re
    # 先处理无序列表
    html = re.sub(r'(<li>.*?</li>\n?)+', lambda m: '<ul>' + m.group(0) + '</ul>', html, flags=re.DOTALL)
    return html

# ==================== /Markdown 转 HTML ====================

# 加载环境变量
env_path = Path.home() / "OneDrive" / ".env"
load_dotenv(env_path, override=True)
# 也尝试加载项目目录下的 .env
load_dotenv(Path(__file__).parent / ".env", override=True)

# 尝试各种可能的 API key 环境变量名
DEEPSEEK_API_KEY = (
    os.getenv("DEEPSEEK_API_KEY") or
    os.getenv("DEEPSEEK_KEY") or
    os.getenv("DEEPSEEK_API") or
    "sk-your-deepseek-api-key-here"
)
DEEPSEEK_BASE_URL = os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com")

import requests

# ==================== Tesseract OCR 配置 ====================
# 自动检测 Tesseract 安装路径
_tesseract_paths = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    os.path.expandvars(r"%LOCALAPPDATA%\Tesseract-OCR\tesseract.exe"),
    os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe"),
]
_TESSERACT_EXE = None
for _tp in _tesseract_paths:
    if os.path.exists(_tp):
        _TESSERACT_EXE = _tp
        break

# 配置 TESSDATA 前缀（语言包目录）
_TESSDATA_DIRS = [
    os.path.expandvars(r"%LOCALAPPDATA%\tesseract-ocr\tessdata"),
]
os.environ.setdefault("TESSDATA_PREFIX", _TESSDATA_DIRS[0])

if _TESSERACT_EXE:
    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_EXE
        print(f"[OCR] Tesseract 已就绪: {_TESSERACT_EXE}")
    except ImportError:
        print("[OCR] pytesseract 未安装，OCR 功能将回退到 API 模式")
else:
    print("[OCR] 未找到 Tesseract 安装，OCR 功能将回退到 API 模式")

# ==================== RapidOCR 配置（中英文通用识别） ====================
_rapid_ocr = None
try:
    from rapidocr_onnxruntime import RapidOCR
    _rapid_ocr = RapidOCR()
    print("[OCR] RapidOCR 已就绪（支持中英文混合识别）")
except Exception as e:
    print(f"[OCR] RapidOCR 不可用: {e}")

# ==================== Flask 初始化 ====================
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "doc-platform-secret-2026")
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB
app.config['UPLOAD_FOLDER'] = str(Path(__file__).parent / 'uploads')
app.config['DATABASE'] = str(Path(__file__).parent / 'database' / 'platform.db')
app.config['SESSION_PERMANENT'] = True
app.config['PERMANENT_SESSION_LIFETIME'] = datetime.timedelta(days=7)
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
CORS(app, supports_credentials=True)

# 未登录用户不能执行写操作（但允许 GET 查看空数据和登录注册）
@app.before_request
def require_auth_for_write():
    if request.path.startswith('/api/') and not request.path.startswith('/api/auth/'):
        if request.method in ('POST', 'PUT', 'DELETE') and 'user_id' not in session:
            return jsonify({"error": "请先登录"}), 401

# 确保目录存在
for p in [app.config['UPLOAD_FOLDER'], str(Path(__file__).parent / 'database')]:
    Path(p).mkdir(parents=True, exist_ok=True)

ALLOWED_EXTENSIONS = {
    'pdf', 'doc', 'docx', 'txt', 'epub', 'mobi',
    'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'tiff',
    'mp3', 'wav', 'ogg', 'flac', 'aac', 'm4a',
    'mp4', 'avi', 'mov', 'mkv', 'webm',
    'xlsx', 'xls', 'csv', 'ppt', 'pptx', 'md', 'json', 'xml', 'html'
}

# ==================== 数据库 ====================
db_lock = Lock()

def get_db():
    if 'db' not in g:
        g.db = sqlite3.connect(app.config['DATABASE'])
        g.db.row_factory = sqlite3.Row
        g.db.execute("PRAGMA journal_mode=WAL")
        g.db.execute("PRAGMA foreign_keys=ON")
    return g.db

@app.teardown_appcontext
def close_db(e=None):
    db = g.pop('db', None)
    if db:
        db.close()

def init_db():
    db = sqlite3.connect(app.config['DATABASE'])
    # 设置时区为本地时间，解决 UTC 时间差8小时问题
    db.execute("PRAGMA journal_mode=WAL")
    # 使用本地时间作为默认时间戳
    import time
    offset = time.timezone if (time.localtime().tm_isdst == 0) else time.altzone
    tz_offset_hours = -offset // 3600
    tz_str = f"{'+' if tz_offset_hours >= 0 else ''}{tz_offset_hours:02d}:00"
    db.execute(f"PRAGMA user_version=1")
    db.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            email TEXT UNIQUE,
            password_hash TEXT NOT NULL,
            verified INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS captcha_codes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            email TEXT NOT NULL,
            code TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            used INTEGER DEFAULT 0
        );

        CREATE TABLE IF NOT EXISTS notes (
            id TEXT PRIMARY KEY,
            title TEXT DEFAULT '无标题笔记',
            content TEXT DEFAULT '',
            plain_text TEXT DEFAULT '',
            tags TEXT DEFAULT '[]',
            folder_id TEXT DEFAULT 'root',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            is_deleted INTEGER DEFAULT 0
        );

        CREATE TABLE IF NOT EXISTS note_history (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            note_id TEXT NOT NULL,
            content TEXT,
            title TEXT,
            version INTEGER DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (note_id) REFERENCES notes(id)
        );

        CREATE TABLE IF NOT EXISTS files (
            id TEXT PRIMARY KEY,
            original_name TEXT NOT NULL,
            stored_name TEXT NOT NULL,
            file_path TEXT NOT NULL,
            file_size INTEGER DEFAULT 0,
            file_type TEXT DEFAULT '',
            mime_type TEXT DEFAULT '',
            folder_id TEXT DEFAULT 'root',
            tags TEXT DEFAULT '[]',
            ocr_text TEXT DEFAULT '',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            is_deleted INTEGER DEFAULT 0
        );

        CREATE TABLE IF NOT EXISTS folders (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            parent_id TEXT DEFAULT 'root',
            folder_type TEXT DEFAULT 'normal',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS translations (
            id TEXT PRIMARY KEY,
            source_text TEXT NOT NULL,
            translated_text TEXT NOT NULL,
            source_lang TEXT DEFAULT 'auto',
            target_lang TEXT DEFAULT 'zh',
            context_type TEXT DEFAULT 'text',
            audio_file_id TEXT,
            tags TEXT DEFAULT '[]',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS ai_chats (
            id TEXT PRIMARY KEY,
            title TEXT DEFAULT '新对话',
            messages TEXT DEFAULT '[]',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS templates (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            category TEXT DEFAULT '通用',
            content TEXT NOT NULL,
            description TEXT DEFAULT '',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );

        CREATE TABLE IF NOT EXISTS config (
            key TEXT PRIMARY KEY,
            value TEXT DEFAULT ''
        );

        -- 默认文件夹
        INSERT OR IGNORE INTO folders (id, name, parent_id, folder_type) VALUES
            ('root', '根目录', '', 'root'),
            ('notes_default', '我的笔记', 'root', 'system'),
            ('translations_default', '翻译记录', 'root', 'system'),
            ('ocr_default', 'OCR识别', 'root', 'system');

        -- 清理已移除功能的残留数据（摘抄记录、我的文件等旧文件夹）
        DELETE FROM folders WHERE name IN ('摘抄记录','知识库','excerpts','knowledge','我的文件','作业','论文')
          OR id IN ('files_default');

        -- 默认模板
        INSERT OR IGNORE INTO templates (id, name, category, content, description) VALUES
            ('tpl_meeting', '会议纪要', '工作', '<h2>会议纪要</h2><p><strong>时间：</strong></p><p><strong>地点：</strong></p><p><strong>参会人员：</strong></p><p><strong>主持人：</strong></p><p><strong>记录人：</strong></p><h3>一、会议议题</h3><p></p><h3>二、讨论内容</h3><p></p><h3>三、决议事项</h3><p></p><h3>四、待办事项</h3><p></p>', '标准会议纪要模板'),
            ('tpl_study', '学习笔记', '学习', '<h2>学习笔记</h2><p><strong>主题：</strong></p><p><strong>日期：</strong></p><p><strong>来源：</strong></p><h3>核心概念</h3><p></p><h3>重点内容</h3><p></p><h3>疑问与思考</h3><p></p><h3>总结</h3><p></p>', '通用学习笔记模板'),
            ('tpl_reading', '读书笔记', '学习', '<h2>读书笔记</h2><p><strong>书名：</strong></p><p><strong>作者：</strong></p><p><strong>阅读日期：</strong></p><h3>书籍概要</h3><p></p><h3>精彩摘录</h3><p></p><h3>个人感悟</h3><p></p><h3>行动清单</h3><p></p>', '读书笔记专用模板'),
            ('tpl_report', '工作报告', '工作', '<h2>工作报告</h2><p><strong>报告人：</strong></p><p><strong>时间范围：</strong></p><h3>一、工作概述</h3><p></p><h3>二、完成情况</h3><p></p><h3>三、问题与挑战</h3><p></p><h3>四、下阶段计划</h3><p></p>', '工作报告模板'),
            ('tpl_daily', '每日记录', '生活', '<h2>每日记录</h2><p><strong>日期：</strong></p><p><strong>天气：</strong></p><p><strong>心情：</strong></p><h3>今日事项</h3><p></p><h3>今日收获</h3><p></p><h3>明日计划</h3><p></p>', '日常记录模板');

        -- 默认配置
        INSERT OR IGNORE INTO config (key, value) VALUES
            ('deepseek_api_key', ''),
            ('deepseek_base_url', 'https://api.deepseek.com'),
            ('default_source_lang', 'zh'),
            ('default_target_lang', 'en');
    """)
    # 迁移：给已有表添加 user_id 字段（如果不存在）
    try: db.execute("ALTER TABLE notes ADD COLUMN user_id INTEGER DEFAULT 1")
    except: pass
    try: db.execute("ALTER TABLE files ADD COLUMN user_id INTEGER DEFAULT 1")
    except: pass
    try: db.execute("ALTER TABLE folders ADD COLUMN user_id INTEGER DEFAULT 1")
    except: pass
    try: db.execute("ALTER TABLE translations ADD COLUMN user_id INTEGER DEFAULT 1")
    except: pass
    try: db.execute("ALTER TABLE ai_chats ADD COLUMN user_id INTEGER DEFAULT 1")
    except: pass
    try: db.execute("ALTER TABLE users ADD COLUMN email TEXT UNIQUE")
    except: pass
    try: db.execute("ALTER TABLE users ADD COLUMN verified INTEGER DEFAULT 0")
    except: pass
    # 给已有笔记/文件设置默认 user_id
    db.execute("UPDATE notes SET user_id = 1 WHERE user_id IS NULL")
    db.execute("UPDATE files SET user_id = 1 WHERE user_id IS NULL")
    db.execute("UPDATE folders SET user_id = 1 WHERE user_id IS NULL")
    db.execute("UPDATE translations SET user_id = 1 WHERE user_id IS NULL")
    db.execute("UPDATE ai_chats SET user_id = 1 WHERE user_id IS NULL")
    db.commit()
    db.close()

init_db()

# ==================== 工具函数 ====================

def generate_id():
    return str(uuid.uuid4())[:12]

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_type(ext):
    ext = ext.lower()
    if ext in ('pdf',):
        return 'pdf'
    elif ext in ('doc', 'docx'):
        return 'word'
    elif ext in ('txt', 'md', 'json', 'xml', 'html', 'csv'):
        return 'text'
    elif ext in ('epub', 'mobi'):
        return 'ebook'
    elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'tiff'):
        return 'image'
    elif ext in ('mp3', 'wav', 'ogg', 'flac', 'aac', 'm4a'):
        return 'audio'
    elif ext in ('mp4', 'avi', 'mov', 'mkv', 'webm'):
        return 'video'
    elif ext in ('xlsx', 'xls'):
        return 'spreadsheet'
    elif ext in ('ppt', 'pptx'):
        return 'presentation'
    return 'other'

def get_api_config():
    db = get_db()
    configs = {}
    for row in db.execute("SELECT key, value FROM config").fetchall():
        configs[row['key']] = row['value']
    api_key = configs.get('deepseek_api_key', '') or DEEPSEEK_API_KEY
    base_url = configs.get('deepseek_base_url', '') or DEEPSEEK_BASE_URL
    return api_key, base_url

def call_deepseek(messages, temperature=0.7, stream=False, max_tokens=4096):
    """调用 DeepSeek API"""
    api_key, base_url = get_api_config()
    if api_key == 'sk-your-deepseek-api-key-here':
        return {"error": "请先在设置中配置 DeepSeek API Key"}
    
    try:
        resp = requests.post(
            f"{base_url.rstrip('/')}/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": "deepseek-chat",
                "messages": messages,
                "temperature": temperature,
                "max_tokens": max_tokens,
                "stream": stream
            },
            timeout=120
        )
        resp.raise_for_status()
        data = resp.json()
        return {"content": data["choices"][0]["message"]["content"]}
    except requests.exceptions.Timeout:
        return {"error": "API 请求超时，请重试"}
    except requests.exceptions.RequestException as e:
        return {"error": f"API 请求失败: {str(e)}"}

def extract_text_from_file(filepath, file_type):
    """从文件中提取文本"""
    text = ""
    try:
        if file_type == 'text':
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            for enc in encodings:
                try:
                    with open(filepath, 'r', encoding=enc) as f:
                        text = f.read()
                    break
                except:
                    continue
        elif file_type == 'pdf':
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(filepath)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
            except:
                pass
        elif file_type == 'word':
            try:
                ext = filepath.rsplit('.', 1)[-1].lower() if '.' in filepath else ''
                if ext == 'docx':
                    from docx import Document
                    doc = Document(filepath)
                    # 提取段落文字
                    for para in doc.paragraphs:
                        if para.text.strip():
                            text += para.text + "\n"
                    # 提取表格文字
                    for table in doc.tables:
                        for row in table.rows:
                            row_texts = []
                            for cell in row.cells:
                                ct = cell.text.strip()
                                if ct:
                                    row_texts.append(ct)
                            if row_texts:
                                text += " | ".join(row_texts) + "\n"
                        text += "\n"
                    # 提取页眉页脚
                    for section in doc.sections:
                        if section.header:
                            for para in section.header.paragraphs:
                                if para.text.strip():
                                    text += "[页眉] " + para.text + "\n"
                        if section.footer:
                            for para in section.footer.paragraphs:
                                if para.text.strip():
                                    text += "[页脚] " + para.text + "\n"
                elif ext == 'doc':
                    # 旧版 .doc 格式，尝试用 antiword 或直接提示
                    text = "[旧版 .doc 格式暂不支持自动提取，请转为 .docx 格式后重新上传]"
            except ImportError:
                text = "[缺少 python-docx 库，请运行: pip install python-docx]"
            except Exception as e:
                import traceback
                print(f"[WARN] Word extract error: {traceback.format_exc()}")
                text = f"[Word文档提取失败: {str(e)}]"


        elif file_type == 'ebook':
            try:
                from ebooklib import epub
                book = epub.read_epub(filepath)
                for item in book.get_items():
                    if item.get_type() == 9:  # ITEM_DOCUMENT
                        content = item.get_content().decode('utf-8', errors='ignore')
                        # 简单去除 HTML 标签
                        content = re.sub('<[^>]+>', ' ', content)
                        text += content + "\n"
            except:
                pass
        elif file_type == 'spreadsheet':
            # Excel 文件 (xlsx/xls)
            try:
                ext = filepath.rsplit('.', 1)[-1].lower() if '.' in filepath else ''
                if ext == 'xlsx':
                    from openpyxl import load_workbook
                    wb = load_workbook(filepath, read_only=True, data_only=True)
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        text += f"[工作表: {sheet_name}]\n"
                        for row in ws.iter_rows(values_only=True):
                            row_data = [str(cell) if cell is not None else '' for cell in row]
                            if any(row_data):
                                text += '\t'.join(row_data) + '\n'
                        text += '\n'
                    wb.close()
                elif ext == 'xls':
                    import xlrd
                    wb = xlrd.open_workbook(filepath)
                    for sheet in wb.sheets():
                        text += f"[工作表: {sheet.name}]\n"
                        for row_idx in range(sheet.nrows):
                            row_data = []
                            for col_idx in range(sheet.ncols):
                                cell_value = sheet.cell_value(row_idx, col_idx)
                                row_data.append(str(cell_value))
                            if any(row_data):
                                text += '\t'.join(row_data) + '\n'
                        text += '\n'
            except Exception as e:
                import traceback
                text = f"[Excel文件提取失败: {str(e)}]"
                print(f"[WARN] Excel extract error: {traceback.format_exc()}")
        elif file_type == 'presentation':
            # PowerPoint 文件 (pptx/ppt)
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                for slide_num, slide in enumerate(prs.slides, 1):
                    text += f"[幻灯片 {slide_num}]\n"
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    text += para.text.strip() + "\n"
                        # 提取表格内容
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_texts = [cell.text.strip() for cell in row.cells]
                                if any(row_texts):
                                    text += " | ".join(row_texts) + "\n"
                    text += "\n"
            except Exception as e:
                import traceback
                text = f"[PPT文件提取失败: {str(e)}]"
                print(f"[WARN] PPT extract error: {traceback.format_exc()}")
        elif file_type == 'csv':
            try:
                with open(filepath, 'r', encoding='utf-8-sig') as f:
                    text = f.read()
            except UnicodeDecodeError:
                try:
                    with open(filepath, 'r', encoding='gbk') as f:
                        text = f.read()
                except:
                    text = "[CSV文件编码不支持]"
    except Exception as e:
        text = f"[提取文本失败: {str(e)}]"
    return text.strip()

# ==================== API 路由 ====================

# ---------- 笔记相关 ----------

@app.route('/api/notes', methods=['GET'])
def api_get_notes():
    db = get_db()
    folder = request.args.get('folder', 'root')
    search = request.args.get('search', '')
    tag = request.args.get('tag', '')
    
    uid = get_current_user_id()
    query = "SELECT * FROM notes WHERE is_deleted = 0 AND user_id = ?"
    params = [uid]
    if folder:
        query += " AND folder_id = ?"
        params.append(folder)
    if search:
        query += " AND (title LIKE ? OR plain_text LIKE ?)"
        params.extend([f'%{search}%', f'%{search}%'])
    if tag:
        query += " AND tags LIKE ?"
        params.append(f'%"{tag}"%')
    query += " ORDER BY updated_at DESC"
    
    notes = []
    for row in db.execute(query, params).fetchall():
        n = dict(row)
        n['tags'] = json.loads(n['tags']) if n['tags'] else []
        notes.append(n)
    return jsonify(notes)

@app.route('/api/notes', methods=['POST'])
def api_create_note():
    data = request.get_json()
    note_id = generate_id()
    title = data.get('title', '无标题笔记')
    content = data.get('content', '')
    tags = json.dumps(data.get('tags', []), ensure_ascii=False)
    folder_id = data.get('folder_id', 'notes_default')
    
    # AI分析报告：将Markdown转为HTML后存储，确保完整渲染
    if title.startswith('[AI分析报告]') or title.startswith('[多文件对比分析]'):
        if not content.startswith('<'):
            print(f"[DEBUG] 转换Markdown→HTML, 原始长度:{len(content)}")
            content = markdown_to_html(content)
            print(f"[DEBUG] 转换后长度:{len(content)}, 开头10字符:{content[:50]}")
        else:
            print(f"[DEBUG] 内容已包含HTML标签，跳过转换")
    
    # 提取纯文本
    plain_text = re.sub('<[^>]+>', ' ', content)
    
    db = get_db()
    db.execute(
        "INSERT INTO notes (id, title, content, plain_text, tags, folder_id, user_id) VALUES (?,?,?,?,?,?,?)",
        (note_id, title, content, plain_text, tags, folder_id, get_current_user_id())
    )
    # 创建初始历史版本
    db.execute(
        "INSERT INTO note_history (note_id, content, title, version) VALUES (?,?,?,1)",
        (note_id, content, title)
    )
    db.commit()
    return jsonify({"id": note_id, "message": "创建成功"})

@app.route('/api/notes/<note_id>', methods=['GET'])
def api_get_note(note_id):
    db = get_db()
    row = db.execute("SELECT * FROM notes WHERE id = ? AND user_id = ?", (note_id, get_current_user_id())).fetchone()
    if not row:
        return jsonify({"error": "笔记不存在"}), 404
    note = dict(row)
    note['tags'] = json.loads(note['tags']) if note['tags'] else []
    return jsonify(note)

@app.route('/api/notes/<note_id>', methods=['PUT'])
def api_update_note(note_id):
    data = request.get_json()
    db = get_db()
    old = db.execute("SELECT content, title FROM notes WHERE id = ? AND user_id = ?", (note_id, get_current_user_id())).fetchone()
    if not old:
        return jsonify({"error": "笔记不存在"}), 404
    
    # 获取当前最新版本号（从 note_history 表）
    ver_row = db.execute(
        "SELECT MAX(version) as max_ver FROM note_history WHERE note_id = ?",
        (note_id,)
    ).fetchone()
    new_version = (ver_row['max_ver'] or 0) + 1
    
    new_title = data.get('title', old['title'])
    new_content = data.get('content', old['content'])
    
    # AI分析报告：将Markdown转为HTML后存储，确保完整渲染
    if new_title.startswith('[AI分析报告]') or new_title.startswith('[多文件对比分析]'):
        if not new_content.startswith('<'):
            new_content = markdown_to_html(new_content)
    
    # 保存历史版本
    db.execute(
        "INSERT INTO note_history (note_id, content, title, version) VALUES (?,?,?,?)",
        (note_id, new_content, new_title, new_version)
    )
    
    plain_text = re.sub('<[^>]+>', ' ', new_content)
    db.execute(
        "UPDATE notes SET title=?, content=?, plain_text=?, tags=?, folder_id=?, updated_at=CURRENT_TIMESTAMP WHERE id=?",
        (
            new_title,
            new_content,
            plain_text,
            json.dumps(data.get('tags', []), ensure_ascii=False),
            data.get('folder_id', 'notes_default'),
            note_id
        )
    )
    db.commit()
    return jsonify({"message": "更新成功"})

@app.route('/api/notes/<note_id>', methods=['DELETE'])
def api_delete_note(note_id):
    db = get_db()
    db.execute("UPDATE notes SET is_deleted = 1 WHERE id = ?", (note_id,))
    db.commit()
    return jsonify({"message": "已删除"})

@app.route('/api/notes/<note_id>/history', methods=['GET'])
def api_note_history(note_id):
    db = get_db()
    rows = db.execute(
        "SELECT id, version, title, created_at FROM note_history WHERE note_id=? ORDER BY version DESC",
        (note_id,)
    ).fetchall()
    return jsonify([dict(r) for r in rows])

@app.route('/api/notes/<note_id>/history/<int:version_id>', methods=['GET'])
def api_get_history_version(note_id, version_id):
    db = get_db()
    row = db.execute(
        "SELECT * FROM note_history WHERE id=?", (version_id,)
    ).fetchone()
    if not row:
        return jsonify({"error": "版本不存在"}), 404
    return jsonify(dict(row))

@app.route('/api/templates', methods=['GET'])
def api_get_templates():
    db = get_db()
    category = request.args.get('category', '')
    query = "SELECT * FROM templates"
    params = []
    if category:
        query += " WHERE category = ?"
        params.append(category)
    rows = db.execute(query, params).fetchall()
    return jsonify([dict(r) for r in rows])

@app.route('/api/templates/<tpl_id>', methods=['GET'])
def api_get_template(tpl_id):
    db = get_db()
    row = db.execute("SELECT * FROM templates WHERE id = ?", (tpl_id,)).fetchone()
    if not row:
        return jsonify({"error": "模板不存在"}), 404
    return jsonify(dict(row))

# ---------- 文件夹 ----------

@app.route('/api/folders', methods=['GET'])
def api_get_folders():
    db = get_db()
    parent = request.args.get('parent', 'root')
    ftype = request.args.get('type', '')
    uid = get_current_user_id()
    query = "SELECT * FROM folders WHERE parent_id = ? AND user_id = ?"
    params = [parent, uid]
    if ftype:
        query += " AND folder_type = ?"
        params.append(ftype)
    rows = db.execute(query, params).fetchall()
    return jsonify([dict(r) for r in rows])

@app.route('/api/folders', methods=['POST'])
def api_create_folder():
    data = request.get_json()
    fid = generate_id()
    db = get_db()
    db.execute(
        "INSERT INTO folders (id, name, parent_id, folder_type, user_id) VALUES (?,?,?,?,?)",
        (fid, data['name'], data.get('parent_id', 'root'), data.get('folder_type', 'normal'), get_current_user_id())
    )
    db.commit()
    return jsonify({"id": fid})

@app.route('/api/folders/<fid>', methods=['DELETE'])
def api_delete_folder(fid):
    db = get_db()
    db.execute("DELETE FROM folders WHERE id = ? AND folder_type != 'system'", (fid,))
    db.commit()
    return jsonify({"message": "已删除"})

# ---------- 文件管理 ----------

@app.route('/api/files', methods=['GET'])
def api_get_files():
    db = get_db()
    folder = request.args.get('folder', 'root')
    ftype = request.args.get('type', '')
    search = request.args.get('search', '')
    
    uid = get_current_user_id()
    query = "SELECT * FROM files WHERE is_deleted = 0 AND folder_id = ? AND user_id = ?"
    params = [folder, uid]
    if ftype:
        query += " AND file_type = ?"
        params.append(ftype)
    if search:
        query += " AND (original_name LIKE ? OR ocr_text LIKE ?)"
        params.extend([f'%{search}%', f'%{search}%'])
    query += " ORDER BY created_at DESC"
    
    rows = db.execute(query, params).fetchall()
    files = []
    for row in rows:
        f = dict(row)
        f['tags'] = json.loads(f['tags']) if f['tags'] else []
        files.append(f)
    return jsonify(files)

@app.route('/api/files/upload', methods=['POST'])
def api_upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "没有文件"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "文件名为空"}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "不支持的文件类型"}), 400
    
    # 保存原始文件名（支持中文），生成安全的存储文件名
    original_filename = file.filename
    # 直接从原始文件名获取扩展名（更可靠）
    ext = original_filename.rsplit('.', 1)[-1].lower() if '.' in original_filename else ''
    if ext not in ALLOWED_EXTENSIONS:
        ext = ''
    filename = secure_filename(file.filename)
    if not filename or ('.' not in filename and ext):
        # secure_filename 把中文名清理掉了，用原始扩展名重建
        filename = f"upload_{generate_id()}.{ext}"
    elif not filename:
        filename = f"upload_{generate_id()}"
        if ext:
            filename += f".{ext}"
    stored_name = f"{generate_id()}.{ext}"
    folder_id = request.form.get('folder_id', 'root')
    
    filepath = Path(app.config['UPLOAD_FOLDER']) / stored_name
    file.save(str(filepath))
    
    file_type = get_file_type(ext)
    file_size = filepath.stat().st_size
    mime_type = file.content_type or 'application/octet-stream'
    
    # 对文本类文件自动提取文字
    ocr_text = ""
    if file_type in ('pdf', 'text', 'word', 'ebook', 'spreadsheet', 'presentation'):
        try:
            ocr_text = extract_text_from_file(str(filepath), file_type)
            print(f"[INFO] 文件 {original_filename} 自动提取文字: {len(ocr_text)} 字符")
        except Exception as e:
            import traceback
            print(f"[WARN] 文件 {original_filename} 提取文字失败: {e}\n{traceback.format_exc()}")
            ocr_text = ""
    
    db = get_db()
    fid = generate_id()
    db.execute(
        """INSERT INTO files (id, original_name, stored_name, file_path, file_size, file_type, mime_type, folder_id, ocr_text, user_id)
           VALUES (?,?,?,?,?,?,?,?,?,?)""",
        (fid, original_filename, stored_name, str(filepath), file_size, file_type, mime_type, folder_id, ocr_text, get_current_user_id())
    )
    db.commit()
    
    return jsonify({
        "id": fid,
        "name": original_filename,
        "type": file_type,
        "size": file_size,
        "ocr_text_preview": ocr_text[:500] if ocr_text else ""
    })

@app.route('/api/files/<fid>/download', methods=['GET'])
def api_download_file(fid):
    db = get_db()
    row = db.execute("SELECT * FROM files WHERE id = ? AND is_deleted = 0 AND user_id = ?", (fid, get_current_user_id())).fetchone()
    if not row:
        return jsonify({"error": "文件不存在"}), 404
    filepath = row['file_path']
    if not os.path.exists(filepath):
        return jsonify({"error": "文件已丢失"}), 404
    return send_file(filepath, as_attachment=True, download_name=row['original_name'])

@app.route('/api/files/<fid>', methods=['DELETE'])
def api_delete_file(fid):
    db = get_db()
    db.execute("UPDATE files SET is_deleted = 1 WHERE id = ? AND user_id = ?", (fid, get_current_user_id()))
    db.commit()
    return jsonify({"message": "已删除"})

@app.route('/api/files/<fid>/ocr', methods=['POST'])
def api_file_ocr(fid):
    """对上传的图片/PDF文件进行 OCR 文字识别"""
    db = get_db()
    row = db.execute("SELECT * FROM files WHERE id = ? AND is_deleted = 0 AND user_id = ?", (fid, get_current_user_id())).fetchone()
    if not row:
        return jsonify({"error": "文件不存在"}), 404
    
    filepath = row['file_path']
    if not os.path.exists(filepath):
        return jsonify({"error": "文件已丢失"}), 404
    
    # 图片类型：优先使用本地 pytesseract OCR，更准确更快速
    if row['file_type'] == 'image':
        return _ocr_with_tesseract(filepath, fid)
    
    # 如果是 PDF，从中提取文本
    elif row['file_type'] == 'pdf':
        text = extract_text_from_file(filepath, 'pdf')
        # 如果 PDF 纯文本提取为空，尝试用 OCR 识别（扫描版 PDF）
        if not text or len(text.strip()) < 20:
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(filepath, first_page=1, last_page=3)
                all_text = []
                for i, img in enumerate(images):
                    try:
                        import pytesseract
                        page_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                        if page_text.strip():
                            all_text.append(page_text)
                    except ImportError:
                        break
                if all_text:
                    text = "\n--- 第{}页 ---\n".format(1).join(
                        f"\n--- 第{i+1}页 ---\n{t}" for i, t in enumerate(all_text)
                    )
            except Exception:
                pass
        db.execute("UPDATE files SET ocr_text = ? WHERE id = ?", (text, fid))
        db.commit()
        return jsonify({"text": text})
    
    # Word / 其他文本类文件：直接用 extract_text_from_file 提取文字
    elif row['file_type'] in ('word', 'text', 'ebook', 'spreadsheet', 'presentation'):
        text = extract_text_from_file(filepath, row['file_type'])
        if text:
            db.execute("UPDATE files SET ocr_text = ? WHERE id = ?", (text.strip(), fid))
            db.commit()
            return jsonify({"text": text.strip()})
    
    return jsonify({"error": "不支持的文件类型"}), 400


def _preprocess_image_for_ocr(img, for_code=False):
    """对图片进行预处理以提升 OCR 准确率"""
    from PIL import ImageOps, ImageFilter
    
    # 转为灰度图
    if img.mode != 'L':
        img = img.convert('L')
    
    # 代码截图：先放大1.5倍提升小字识别率（在二值化之前）
    if for_code:
        w, h = img.size
        img = img.resize((int(w * 1.5), int(h * 1.5)))
    
    # 自动对比度增强
    img = ImageOps.autocontrast(img, cutoff=1)
    # 轻微锐化
    img = img.filter(ImageFilter.SHARPEN)
    
    # 自适应二值化
    pixels = list(img.getdata())
    threshold = sum(pixels) // len(pixels)
    if for_code:
        threshold = int(threshold * 0.75)
    else:
        threshold = int(threshold * 0.85)
    img = img.point(lambda x: 0 if x < threshold else 255, '1')
    img = img.convert('L')
    
    return img


def _ocr_with_tesseract(filepath, fid):
    """使用 OCR 进行文字识别：RapidOCR（中英文通用）优先，Tesseract（代码）兜底"""
    db = get_db()
    from PIL import Image

    img = Image.open(filepath)

    # ===== 第一步：用 RapidOCR 做通用文字识别（中文+英文） =====
    if _rapid_ocr is not None:
        try:
            ocr_result = _rapid_ocr(np.array(img))
            result = ocr_result[0]
            if result and len(result) > 0:
                raw_text = "\n".join([line[1] for line in result])
                elapsed_total = sum(ocr_result[1]) if isinstance(ocr_result[1], list) else 0
                print(f"[OCR] RapidOCR OK, elapsed {elapsed_total:.2f}s, chars: {len(raw_text)}")

                if _is_code_text(raw_text):
                    text = _fix_code_ocr(raw_text)
                    note = "代码识别(RapidOCR)"
                else:
                    text = _postprocess_ocr_text(raw_text)
                    note = "通用文字识别(RapidOCR)"

                if text and len(text.strip()) > 3:
                    db.execute("UPDATE files SET ocr_text = ? WHERE id = ? ", (text.strip(), fid))
                    db.commit()
                    return jsonify({"text": text.strip(), "method": "rapidocr", "note": note})
        except Exception as e:
            print(f"[OCR] RapidOCR 异常: {e}")
            import traceback
            traceback.print_exc()

    # ===== 第二步：回退到 Tesseract（主要用于代码） =====
    try:
        import pytesseract

        # 先尝试英文模式检测代码
        proc_img = _preprocess_image_for_ocr(img, for_code=True)
        raw_text = pytesseract.image_to_string(proc_img, lang='eng', config=r'--oem 3 --psm 3')

        if _is_code_text(raw_text):
            text = _fix_code_ocr(raw_text)
            note = "代码识别(Tesseract)"
        else:
            # 非代码：尝试中英文混合模式
            proc_img2 = _preprocess_image_for_ocr(img, for_code=False)
            raw_text2 = pytesseract.image_to_string(proc_img2, lang='chi_sim+eng', config=r'--oem 3 --psm 6')
            if len(raw_text2.strip()) > len(raw_text.strip()):
                raw_text = raw_text2
            text = _postprocess_ocr_text(raw_text)
            note = "通用文字识别(Tesseract)"

        if text and len(text.strip()) > 5:
            db.execute("UPDATE files SET ocr_text = ? WHERE id = ?", (text.strip(), fid))
            db.commit()
            return jsonify({"text": text.strip(), "method": "tesseract", "note": note})
    except ImportError:
        print("[OCR] pytesseract 未安装")
    except Exception as e:
        import traceback
        print(f"[OCR] Tesseract识别异常: {e}")
        traceback.print_exc()

    # ===== 第三步：最终回退 —— DeepSeek Vision API =====
    return _ocr_with_deepseek_vision(filepath, fid)


def _is_code_text(text):
    """检测OCR文本是否为代码"""
    if not text:
        return False
    import re
    lines = [l for l in text.split('\n') if l.strip()]
    # == 关键修复：移除 len(lines) < 2 早期返回 ==
    # OCR经常将代码压成一行（如 plt.scatter(data['a'], data['b'], c=data['c'])）
    # 多行要求会误判大量单行代码。
    # 改用其他启发式特征（关键字、特殊符号密度、方法调用等）综合判断。

    # 代码关键词（正常格式）
    code_keywords = ['def ', 'if __name__', 'import ', 'from ', 'class ',
                     'for ', 'while ', 'elif ', 'else:', 'try:', 'except',
                     'return ', 'lambda ', 'yield ', 'raise ']
    kw_count = sum(1 for kw in code_keywords if kw in text)

    # 粘连关键词检测（OCR常见问题：空格丢失导致关键字粘连）
    glued_patterns = [
        r'\bfor\w*(?:in|range)\b',          # for...in, for...range (粘连)
        r'\b(?:if|elif|while)\s*\w+\s*:',   # if/elif/while 后跟冒号（可能粘连）
        r'\bdef\s*\w+\s*\(',                 # def funcname(
        r'\bimport\s+\w+',                   # import module
        r'\bfrom\s+\w+\s+import',            # from ... import
        r'\bclass\s+\w+\s*[:\(]',            # class ClassName
        r'\breturn\b.*[=<>!]',               # return with operators
        r'(?<!\w)[{}\[\]()](?!\w)',         # 括号/花括号密集出现
        r'\bprint\s*\(',                     # print(
        r'\bCounter\(\)',                    # Counter()
        r'#Step\d|#step\d',                  # 注释标记
        r'\bpair_counter\b',                # 典型变量名模式
        r'\bfrequent_\w+\b',                 # frequent_items/pairs 等
        r'\bsorted\s*\(',                    # sorted(
        r'\bcombinations?\s*\(',             # combinations(
        r'\bmin_support_count\b',           # 变量名
        r'\bfiltered_transaction\b',        # 长变量名
        r'\btransaction[s]?\b',             # transaction(s)
    ]
    glued_count = sum(1 for p in glued_patterns if re.search(p, text))

    # 行首有缩进（空格/制表符开头）
    indented_lines = sum(1 for l in lines if l[0] in (' ', '\t'))
    indent_ratio = indented_lines / max(len(lines), 1)
    # 特殊符号密度
    special_chars = sum(1 for c in text if c in '(){}[]=<>+-*/%&|^~!:;@#$')
    total_visible = len(text.replace(' ', '').replace('\n', ''))
    special_ratio = special_chars / max(total_visible, 1)

    # == 新增：代码特有但无关键词的模式 ==
    # 适用于 OCR 把代码压成一行但无 import/def/for 等关键词的场景
    # 例如: plt.scatter(data['x'], data['y'], c=data['z'])
    #       x = 10 y = 20 z = x + y print(z) # result
    #       # Plot clusters plt.scatter(...
    code_signals = 0
    # (a) 点号方法调用: .method( 如 .scatter(
    if re.search(r'\.\w+\s*\(', text):
        code_signals += 2
    # (b) 井号注释
    if re.search(r'#\s*\w', text):
        code_signals += 2
    # (c) 下标索引: var['xxx'] 或 var["xxx"]
    if re.search(r"\w+\s*\[['\"]", text):
        code_signals += 2
    # (d) 关键字参数: name=value 形式 (在括号内)
    if re.search(r'\([^)]*,\s*\w+\s*=', text):
        code_signals += 2
    # (e) 显式赋值链: x = ... y = ... z = ... (多个赋值在一行)
    assignments = re.findall(r'(?<!\w)(\w+)\s*=\s*', text)
    if len(set(assignments)) >= 2:
        code_signals += 2

    # 综合判断
    has_normal_keywords = kw_count >= 2
    has_glued_code = glued_count >= 3
    has_glued_code_single = glued_count >= 2 and len(lines) <= 2  # 单行宽松
    has_indent = indent_ratio > 0.3
    high_special = special_ratio > 0.15
    has_code_signals = code_signals >= 3

    return (has_normal_keywords or has_glued_code or has_glued_code_single or
            has_indent or high_special or has_code_signals)


def _split_glued_code_words(text):
    """拆分 OCR 粘连的代码词（空格丢失导致的关键字/标识符粘连）

    例如:
      fortransactionin  → for transaction in
      #Step2:Generateonlypairsfromfrequentitems  → #Step2: Generate only pairs from frequent items
      ifcount >=  → if count >=
      defcalc(     → def calc(
    """
    import re

    # Python 关键字集合（用于识别粘连词中的关键字部分）
    py_keywords = set([
        'and', 'as', 'assert', 'async', 'await', 'break', 'class', 'continue',
        'def', 'del', 'elif', 'else', 'except', 'finally', 'for', 'from',
        'global', 'if', 'import', 'in', 'is', 'lambda', 'nonlocal', 'not',
        'or', 'pass', 'raise', 'return', 'try', 'while', 'with', 'yield',
        'True', 'False', 'None'
    ])

    # 常见内置函数/类型名（也常被粘连）
    builtins = set([
        'print', 'len', 'range', 'str', 'int', 'float', 'list', 'dict', 'set',
        'tuple', 'sorted', 'map', 'filter', 'zip', 'enumerate', 'min', 'max',
        'sum', 'abs', 'open', 'type', 'isinstance', 'issubclass', 'hasattr',
        'getattr', 'setattr', 'input', 'Counter', 'combinations', 'items',
        'keys', 'values', 'append', 'extend', 'insert', 'remove', 'pop',
        'count', 'index', 'find', 'replace', 'strip', 'split', 'join',
        'lower', 'upper', 'start', 'end', 'format'
    ])

    lines = text.split('\n')
    result_lines = []

    for line in lines:
        original = line
        stripped = line.strip()
        if not stripped:
            result_lines.append(line)
            continue

        # 跳过纯注释行（不做拆分）
        if stripped.startswith('#') and not stripped.startswith('#Step'):
            # 但仍处理注释中的粘连
            pass

        # === 规则1：拆分行首粘连关键字 ===
        # forxxx → for xxx, ifxxx → if xxx, whilexxx → while xxx, etc.
        def _split_prefix(m):
            kw = m.group(1)
            rest = m.group(2)
            # 如果rest以大写字母开头，说明是驼峰变量名的一部分，不拆分
            if rest and rest[0].isupper() and len(rest) <= 20 and '_' not in rest[:5]:
                return m.group(0)
            return kw + ' ' + rest

        line = re.sub(r'^(?=\s*)(for|if|elif|while|def|class|from|import|return|yield|raise|with|try|except|finally|else|and|or|not|in|is|print|assert)([a-zA-Z_])',
                      _split_prefix, line)

        # === 规则2：#Step 后面紧跟的粘连文本 ===
        # #Step2:Generateonlypairsfromfrequentitems → #Step2: Generate only pairs...
        def _split_step_comment(m):
            # 【关键修复】group(1)是num，group(2)是space，group(3)才是内容！
            content = m.group(3) if len(m.groups()) >= 3 else m.group(1)
            # 在大写字母前插入空格（拆分驼峰/PascalCase单词）
            split = re.sub(r'([a-z])([A-Z])', r'\1 \2', content)
            # 再在关键词边界处拆分常见英文小词
            small_words = ('from', 'only', 'pairs', 'with', 'that', 'this', 'into', 'over', 'upon',
                           'when', 'then', 'than', 'also', 'some', 'such', 'just', 'more')
            for sw in small_words:
                split = re.sub(r'(?i)(' + sw + r')([a-z])', r'\1 \2', split)
                split = re.sub(r'(?i)([a-z])(' + sw + r')', r'\1 \2', split)
            return '#Step' + m.group('num') + ':' + (m.group('space') or '') + ' ' + split

        line = re.sub(r'#Step(?P<num>\d+):(?P<space>\s*)?(.+)',
                      _split_step_comment, line)

        # 【修复】如果这行已经被规则2处理为完整的 #Step 注释，跳过后续的通用粘连拆分规则
        # 防止规则3.5把注释内容再次错误处理（如把 "from frequent" 再次破坏）
        _is_step_comment_processed = bool(re.match(r'#\s*Step\d+\s*:', stripped))

        # === 规则3：关键字+标识符粘连（不在行首的情况） ===
        # 例如: x=foritem → x = for item, )fortransaction → ) for transaction
        # 【跳过】已处理的#Step注释行不需要此步骤
        # 先定义辅助函数（无论是否跳过都需要，避免 UnboundLocalError）
        def _split_long_glued_chain(match):
            word = match.group(0)
            if len(word) < 10:
                return word

            step1 = re.sub(r'([a-z])([A-Z])', r'\1 \2', word)

            def _split_by_keywords(segment):
                if len(segment) < 6:
                    return segment

                all_kws_lower = {kw.lower() for kw in (py_keywords | builtins)}
                exclude = {'in', 'is', 'on', 'or', 'at', 'it', 'as', 'to',
                           'do', 'go', 'no', 'of', 'up', 'my', 'an', 'be',
                           'er', 'count', 'items'}
                search_kws = all_kws_lower - exclude
                if not search_kws:
                    return segment

                seg_lower = segment.lower()
                kw_matches = []

                for kw in sorted(search_kws, key=len, reverse=True):
                    start = 0
                    while True:
                        pos = seg_lower.find(kw, start)
                        if pos == -1:
                            break
                        kw_matches.append((pos, pos + len(kw), kw))
                        start = pos + 1

                if not kw_matches:
                    return segment

                if len(kw_matches) < 2:
                    total_covered = sum(e - s for s, e, k in kw_matches)
                    is_very_long = len(segment) > 15
                    if not is_very_long and total_covered < len(segment) * 0.5:
                        return segment

                kw_matches.sort(key=lambda x: (x[0], -(x[1] - x[0])))
                filtered = []
                for m in kw_matches:
                    if filtered and m[0] < filtered[-1][1]:
                        continue
                    filtered.append(m)
                kw_matches = filtered

                if not kw_matches:
                    return segment

                parts = []
                last_end = 0
                for start, end, kw in kw_matches:
                    if start > last_end:
                        parts.append(segment[last_end:start])
                    parts.append(segment[start:end])
                    last_end = end
                if last_end < len(segment):
                    parts.append(segment[last_end:])

                return ' '.join(parts)

            segments = step1.split(' ')
            result_parts = []
            for seg in segments:
                result_parts.append(_split_by_keywords(seg))

            return ' '.join(result_parts)

        # 【跳过】已处理的#Step注释行不需要规则3/3.5/4（但函数已定义供后续使用）
        if not _is_step_comment_processed:

            def _split_kw_after_punct(m):
                punct = m.group(1)
                kw = m.group(2)
                rest = m.group(3)
                return punct + ' ' + kw + ' ' + rest

            # 标点后紧接关键字的模式
            line = re.sub(
                r'([:;=,\[\(])(for|if|while|in|and|or|not|is|return|print|from|import|def)([a-zA-Z_])',
                _split_kw_after_punct, line
            )

        # 【修复】更严格的正则：只匹配关键字开头+长尾(6+) 或 超长纯词(12+)
        # 避免误伤 counter(7)/filtered(9)/frequent(9) 等正常标识符
        line = re.sub(
            r'\b(?:for|from|while|if|elif|def|class|return|import|print|with|try|except)'
            r'[a-z]{6,}\b'
            r'|'
            r'\b[a-z]{12,}\b',
            _split_long_glued_chain, line
        )

        # 迭代拆分：对仍然包含粘连的行再次处理（最多3轮）
        for _ in range(3):
            new_line = re.sub(
                r'\b(for|if|elif|while|from|import|def|class|return|print|and|or|not|in|is)'
                r'([a-z]{3,})'
                r'(in|import|count|items|pairs|trans)\b',
                lambda m: m.group(1) + ' ' + m.group(2) + ' ' + m.group(3),
                line)
            if new_line == line:
                break
            line = new_line

            # === 规则4：通用粘连拆分（基于大小写变化）===
            # filtered_transaction=[item → filtered_transaction = [item （在=前加空格）

            # 在赋值运算符前确保有空格
            line = re.sub(r'(?<=[a-zA-Z0-9_])=', ' = ', line)
            line = re.sub(r'= (?==)', '=', line)  # 恢复 == 为双等号

        # if not _is_step_comment_processed 结束
        result_lines.append(line)

    text = '\n'.join(result_lines)

    # === 全局规则：清理多余空格 ===
    # 多个空格合并为一个（保留缩进）
    cleaned_lines = []
    for line in text.split('\n'):
        if line and line[0] in (' ', '\t'):
            # 保留行首缩进
            indent_match = re.match(r'^(\s+)', line)
            if indent_match:
                indent = indent_match.group(1)
                rest = line[indent_match.end():]
                rest = re.sub(r'  +', ' ', rest).strip()
                cleaned_lines.append(indent + rest)
            else:
                cleaned_lines.append(re.sub(r'  +', ' ', line).rstrip())
        else:
            cleaned_lines.append(re.sub(r'  +', ' ', line).rstrip())

    return '\n'.join(cleaned_lines)


def _fix_code_ocr(text):
    """修复OCR代码：全角转半角 + 拆分粘连词 + 括号平衡 + 智能缩进恢复 + AI兜底"""
    import re

    # 1. 全角标点->半角（完整映射表）
    punct_map = {
        # 基础标点
        '\uff0c': ',', '\uff1b': ';', '\uff1a': ':', '\uff0e': '.', '\uff1f': '?', '\uff01': '!',
        # 括号（关键！OCR代码常见错误）
        '\uff08': '(', '\uff09': ')',
        '\u3008': '<', '\u3009': '>',
        '\u300a': '[', '\u300b': ']',
        '\u300c': '[', '\u300d': ']',
        '\u300e': '[', '\u300f': ']',
        '\u3010': '[', '\u3011': ']',
        '\uff3b': '[', '\uff3d': ']',
        '\u3001': ',',
        # 引号
        '\u201c': '"', '\u201d': '"',
        '\u2018': "'", "\u2019": "'",
        '\uff02': '"',
        '\uff07': "'",
        # 空格和特殊字符
        '\u3000': ' ',
        '\uff05': '%', '\uff06': '&',
        '\uff0a': '*', '\uff0b': '+',
        '\uff0d': '-', '\uff0f': '/',
        '\uff1c': '<', '\uff1e': '>',
        '\uff1d': '=', '\uff3d': '\\',
    }
    for k, v in punct_map.items():
        text = text.replace(k, v)

    # 2. 【核心新增】拆分粘连的标识符/关键字
    # OCR 经常丢失空格导致: fortransaction → for transaction, ifcount → if count
    text = _split_glued_code_words(text)

    # 3. 常见OCR代码错误修正（在ast解析之前尽可能修复）
    ocr_fixes = [
        (r'(?<!\w)0(?=\s*[;)=\]])', 'O'),   # 单独的0可能是O
        (r'(?<=\w)ll(?=\s*\()', '11'),        # ll( 可能是 11(
        (r'\|\|', ' or '),                   # || → or
        (r'&&', ' and '),                    # && → and
        (r'!=+', '!='),                      # 多个= → !=
        (r'==+', '=='),                      # 多个= → ==
        (r"\[([^\]]*)\]\s*$", lambda m: '[' + m.group(1) + ']'),  # 末尾孤立的]
        (r"^(\s*)\}(.*)$", r'\1}\2'),        # 确保独立}行
        # 【新增】模块名中的连字符 → 点号：matplotlib-pyplot → matplotlib.pyplot
        # 规则：两个标识符之间的连字符（两侧无空格），且右侧以小写字母开头
        (r'(?<=[a-zA-Z0-9_])-(?=[a-z][a-zA-Z0-9_]*)', '.'),
    ]
    for pattern, replacement in ocr_fixes:
        text = re.sub(pattern, replacement, text)

    # 3.5 【关键修复】冒号后换行拆分 —— OCR常把两行代码压成一行
    # 例如: "for pair in combinations(..., 2): pair_counter[pair] += 1"
    #       → 拆分为两行，让后续缩进算法能正确处理
    def _split_colon_line(line):
        """将包含冒号+后续语句的单行拆分为多行"""
        stripped = line.strip()
        if not stripped:
            return [line]

        # 快速检查：行中是否有冒号且冒号后有内容（排除字典/切片/类型注解）
        # 策略：找到不在括号/引号内的冒号，如果后面还有代码则拆分

        result_lines = []
        current = ''
        paren_depth = 0
        bracket_depth = 0
        brace_depth = 0
        in_string = None  # None, '"', "'"
        i = 0
        while i < len(stripped):
            ch = stripped[i]

            # 跟踪字符串状态
            if in_string:
                current += ch
                if ch == in_string:
                    # 检查是否是转义引号
                    # 简单处理：假设没有复杂转义
                    in_string = None
                i += 1
                continue

            if ch in ('"', "'"):
                in_string = ch
                current += ch
                i += 1
                continue

            # 跟踪括号深度
            if ch == '(':
                paren_depth += 1
            elif ch == ')':
                paren_depth -= 1
            elif ch == '[':
                bracket_depth += 1
            elif ch == ']':
                bracket_depth -= 1
            elif ch == '{':
                brace_depth += 1
            elif ch == '}':
                brace_depth -= 1

            current += ch

            # 找到冒号：如果在括号外、字符串外，且后面还有非空内容 → 可能需要拆分
            if (ch == ':' and
                paren_depth == 0 and bracket_depth == 0 and brace_depth == 0 and
                not in_string and
                i + 1 < len(stripped) and stripped[i + 1] not in (' ', '\t', '/', '#')):

                rest_of_line = stripped[i + 1:].strip()
                # 只有当冒号后的内容看起来像独立语句时才拆分
                # 排除：切片 [:n]、字典值 {k:v}、类型注解 def f(x: int)
                should_split = True

                # 不拆分的情况：
                # 1. 冒号前是数字或 ] 或 ) （可能是切片）
                before_colon = current.rstrip()[:-1].rstrip() if len(current.rstrip()) > 1 else ''
                if before_colon and before_colon[-1] in (']', ')'):
                    should_split = False
                # 2. 冒号紧跟空格然后是数字/引号/大括号（可能是字典/类型注解）
                if re.match(r'^[\d"\'\{\[]', rest_of_line):
                    should_split = False
                # 3. 当前行以常见不产生块的关键字开头（如 case、default）
                line_start = stripped.lstrip()[:10]
                if any(line_start.startswith(kw) for kw in ('case ', 'http')):
                    should_split = False

                if should_split and rest_of_line:
                    # 将冒号后的部分作为新行
                    result_lines.append(current)
                    current = '    ' + rest_of_line  # 给新行加基础缩进
                    # 继续处理剩余部分（可能还有更多冒号）
                    remaining = '    ' + rest_of_line
                    sub_lines = _split_colon_line(remaining)
                    if len(sub_lines) > 1 or sub_lines[0].strip() != remaining.strip():
                        result_lines.extend(sub_lines[:-1])
                        current = sub_lines[-1] if sub_lines else ''
                    break

            i += 1

        if current:
            result_lines.append(current)

        return result_lines if len(result_lines) > 1 else [line]

    # 对每行应用冒号拆分
    expanded_lines = []
    for ll in text.split('\n'):
        expanded_lines.extend(_split_colon_line(ll))
    text = '\n'.join(expanded_lines)

    # === 3.6 【关键修复】语句边界拆分 —— 处理OCR把多个独立语句粘到一行的情况 ===
    # 典型场景：
    #   - "if count >= min_support_count print("\nFrequent pairs:")"
    #     → if条件是字典/推导式内的过滤器，print是独立语句
    #   - "} return result"  → 两个独立语句
    #   - "items() if x > 0 else y print('done')"  → 表达式 + print
    #
    # 策略：检测行内出现"顶级语句起始关键字"的位置，在该位置切分，
    #       并自动补全前面未闭合的括号/花括号

    # 常见的顶级/语句起始模式（出现在行中间时大概率是新语句开始）
    # 注意：不包含 'if' —— 因为推导式/字典内的 "if filter" 不应触发拆分
    _stmt_start_patterns = [
        (r'\bprint\s*\(', 'print'),       # print( 几乎总是新语句（除非在参数内部，但那种情况罕见）
        (r'\breturn\s+', 'return'),
        (r'\bimport\s+', 'import'),
        (r'\bfrom\s+\w+(\.\w+)*\s+import', 'from import'),  # from X[.Y] import（支持多级模块）
        (r'\bdef\s+', 'def'),
        (r'\bclass\s+', 'class'),
        (r'\belif\s+', 'elif'),
        (r'\belse\s*:', 'else'),
        (r'\bfor\s+', 'for'),
        (r'\bwhile\s+', 'while'),
        (r'\btry\s*:', 'try'),
        (r'\bexcept\s+', 'except'),
        (r'\bfinally\s*:', 'finally'),
        (r'\bwith\s+', 'with'),
        (r'braise\s+', 'raise'),
        (r'\byield\s+', 'yield'),
        (r'\bassert\s+', 'assert'),
        (r'\bpass\b', 'pass'),
        (r'\bbreak\b', 'break'),
        (r'\bcontinue\b', 'continue'),
    ]

    def _split_statement_boundaries(text):
        """在语句边界处拆分行，跟踪跨行括号深度以正确补全闭合括号"""
        lines = text.split('\n')
        result = []

        # 跟踪跨行的累积括号深度（用于在拆分时正确补全闭合括号）
        cumul_depth = {'{': 0, '(': 0, '[': 0}
        br_pairs = {'{': '}', '(': ')', '[': ']'}

        def _update_depth(s, depths):
            """更新累积括号深度"""
            for ch in s:
                if ch in depths:
                    depths[ch] += 1
                elif ch in br_pairs:
                    close_for = {v: k for k, v in br_pairs.items()}
                    if ch in close_for:
                        open_br = close_for[ch]
                        if depths[open_br] > 0:
                            depths[open_br] -= 1

        for line in lines:
            stripped = line.strip()
            if not stripped:
                result.append(line)
                # 空行也更新括号状态（处理跨空行的括号）
                _update_depth(stripped, cumul_depth)
                continue

            # 先更新这一行的括号到累积深度（用于后续行）
            # 但我们需要在拆分前知道"拆分点之前"的括号状态

            # == 预处理：注释后紧跟代码（OCR合并多行导致）==
            # 例如: "# Plot clusters plt.scatter(data['x'], data['y'])"
            # 在Python中#到行尾都是注释，代码被吞掉了！
            # → 应在第一个代码特征前拆分，保留注释行，后续代码另起一行
            best_split_pos = -1
            best_match_name = None
            if stripped.startswith('#') and len(stripped.strip('# \t')) > 5:
                code_after = re.search(r'(\.\w+\s*\(|\w+\s*\[[\'\"])', stripped[1:])
                if code_after:
                    candidate = 1 + code_after.start()
                    if 2 < candidate < len(stripped) - 2:
                        best_split_pos = candidate
                        best_match_name = '#comment'

            # 只在注释预处理没有找到拆分点时，才执行语句起始关键字扫描
            if best_split_pos == -1:
                for pattern, name in _stmt_start_patterns:
                    for m in re.finditer(pattern, stripped):
                        pos = m.start()

                        # 跳过行首匹配（已经是行开头了，不需要拆分）
                        if pos <= 2:
                            continue

                        # 检查这个位置是否在字符串内部
                        before = stripped[:pos]
                        q_before = before.count('"') + before.count("'")
                        if q_before % 2 != 0:
                            continue  # 在字符串内部，跳过

                        # 【关键修复】检查此位置是否在当前行内打开的括号内部
                        # 注意：只用当前行的局部深度（不含跨行累积），因为跨行的未闭合括号
                        # 通常意味着前面语句未结束，而新语句开始时应该可以拆分
                        local_depth = {'{': 0, '(': 0, '[': 0}
                        _update_depth(stripped[:pos], local_depth)
                        total_local_depth = sum(local_depth.values())
                        if total_local_depth > 0:
                            continue  # 在本行打开的括号内，跳过

                        # 【新增】排除 from ... import 中的 import（不是语句边界）
                        # 例如 "from sklearn.cluster import KMeans" 中的 import 不应拆分
                        if name == 'import':
                            before_text = stripped[:pos]
                            # 检查前面是否有 from 关键字 + 模块路径（说明是 from X import Y 的一部分）
                            if re.search(r'\bfrom\s+[\w.]+\s*$', before_text):
                                continue

                        if best_split_pos == -1 or pos < best_split_pos:
                            best_split_pos = pos
                            best_match_name = name

            if best_split_pos > 0:
                # 在此处拆分
                part1 = stripped[:best_split_pos].rstrip()
                part2 = stripped[best_split_pos:].lstrip()

                # 计算part1的括号（只用当前行局部深度，不跨行）
                local_depth = {'{': 0, '(': 0, '[': 0}
                _update_depth(part1, local_depth)

                # 补全part1中未闭合的括号（仅当前行内打开的）
                for br_open in ('{', '(', '['):
                    diff = local_depth[br_open]
                    if diff > 0:
                        part1 += br_pairs[br_open] * diff

                if part1:
                    result.append(part1)

                # 重置累积深度（因为已补全了闭合括号）
                cumul_depth = {'{': 0, '(': 0, '[': 0}

                if part2:
                    result.append(part2)
                    # 更新累积深度为part2的状态
                    _update_depth(part2, cumul_depth)
            else:
                result.append(line)
                # 更新累积深度
                _update_depth(stripped, cumul_depth)

        return '\n'.join(result)

    # 3.6 【关键修复】语句边界拆分 —— 迭代模式确保多语句行完全拆分
    # 单次调用只能拆一个位置，但OCR一行可能有3+个独立语句（如多个import）
    for _ in range(10):  # 最多10轮，防止无限循环
        new_text = _split_statement_boundaries(text)
        if new_text == text:
            break
        text = new_text

    # 4. 保留原始行
    lines = [l.rstrip() for l in text.split('\n')]
    code = '\n'.join(lines)

    # 4. 括号平衡：补全缺失的闭括号
    for br_open, br_close in [('{', '}'), ('(', ')'), ('[', ']')]:
        diff = code.count(br_open) - code.count(br_close)
        if diff > 0:
            code += br_close * diff

    # 5. 尝试用 ast.unparse 格式化（最可靠——Python自己理解代码结构）
    import ast
    try:
        tree = ast.parse(code)
        if hasattr(ast, 'unparse'):
            result = ast.unparse(tree)
            # 验证结果确实包含缩进（说明ast成功了）
            if any(l.startswith(' ') for l in result.split('\n') if l.strip()):
                return result
    except (SyntaxError, IndentationError, ValueError):
        pass

    # 6. ast失败：用改进的堆栈算法推断缩进（关键bug修复：入栈cur+1而非cur）
    result_lines = []
    indent_stack = []  # 记录每个块开始后的缩进级别

    for line in lines:
        s = line.strip()
        if not s:
            result_lines.append('')
            continue

        # 获取当前应有的缩进级别
        cur = indent_stack[-1] if indent_stack else 0

        # elif/else/except/finally：回到所在块的级别（即栈顶）
        if any(s.startswith(k) for k in ('elif ', 'else:', 'except ', 'finally:')):
            if len(indent_stack) >= 1:
                cur = indent_stack[-1]
            else:
                cur = 0

        # 顶级/语句结束符：清空缩进栈回到根
        top_kws = (
            '# ', 'import ', 'from ', 'return ', 'pass', 'break',
            'continue', 'raise ', 'assert ', 'del ', 'yield ',
            'global ', 'nonlocal ',
        )
        # print( 也作为可能的顶级语句（但不在赋值右边的）
        if any(s == kw or s.startswith(kw) for kw in top_kws if kw != 'print('):
            indent_stack.clear()
            cur = 0
        # 特殊处理：print(...) 如果栈很深可能仍是内部调用，保守起见不清栈
        # 但如果print以大写P开头或在行首且栈不太深，考虑清栈

        result_lines.append('    ' * cur + s)

        # 以单个冒号结尾的行是块开始 → 下一行要加缩进
        # 【关键修复】：入栈 cur + 1（下一级缩进），不是 cur！
        if s.endswith(':') and s.count(':') == 1 and not s.startswith('#'):
            indent_stack.append(cur + 1)

    raw_result = '\n'.join(result_lines)

    # 7. 再次尝试 ast 解析（经过缩进修复后可能成功了）
    try:
        tree = ast.parse(raw_result)
        if hasattr(ast, 'unparse'):
            return ast.unparse(tree)
    except (SyntaxError, IndentationError, ValueError):
        pass

    # 8. 终极兜底：调用 DeepSeek AI 修复代码格式
    ai_fixed = _ai_fix_code_format(text)
    if ai_fixed:
        return ai_fixed

    return raw_result


def _ai_fix_code_format(raw_text):
    """调用 DeepSeek AI 修复 OCR 代码的缩进和格式"""
    if not raw_text.strip():
        return None
    try:
        result = call_deepseek([
            {
                "role": "system",
                "content": (
                    "你是一位Python代码修复专家。下面一段Python代码来自OCR（光学字符识别），"
                    "存在以下典型问题：\n"
                    "1. 【缩进丢失】所有行都没有缩进，需要根据代码逻辑恢复4空格缩进\n"
                    "2. 【换行丢失】OCR可能把多行代码合并成一行，例如：\n"
                    "   - 'for x in y: statement' 应拆成两行\n"
                    "   - 'if condition: stmt' 应拆成两行\n"
                    "3. 【粘连词】关键字之间空格丢失，例如：\n"
                    "   - 'fortransactionintransactions' → 'for transaction in transactions'\n"
                    "   - 'fromitertoolsimportcombinations' → 'from itertools import combinations'\n"
                    "   - '#Step2:Generateonlypairs...' → '#Step2: Generate only pairs...'\n"
                    "4. 【括号/标点错误】全角半角混用、括号不匹配等\n\n"
                    "请按以下要求修复：\n"
                    "- 恢复正确的Python语法和4空格缩进\n"
                    "- 把被OCR合并的行正确拆分（特别注意冒号后的语句应在新行）\n"
                    "- 修复粘连的关键字和标识符\n"
                    "- 保持原始代码的逻辑不变\n"
                    "- 只输出修复后的完整Python代码，不要任何解释文字，不要用markdown包裹"
                )
            },
            {
                "role": "user",
                "content": f"请修复以下OCR识别出的Python代码：\n\n{raw_text}"
            }
        ], temperature=0.05, max_tokens=4096)
        if "error" not in result and result.get("content", "").strip():
            content = result["content"].strip()
            # 去掉可能的markdown代码块标记
            if content.startswith("```"):
                # 移除开头的 ```
                lines = content.split("\n")
                # 找到第一个 ``` 后面的行作为开始
                start = 0
                end = len(lines)
                for i, l in enumerate(lines):
                    if l.startswith("```"):
                        if start == i:
                            start = i + 1
                        else:
                            end = i
                            break
                content = "\n".join(lines[start:end]).strip()
            return content
    except Exception as e:
        print(f"[WARN] AI代码修复失败: {e}")
    return None


def _postprocess_ocr_text(text):
    """对 OCR 结果进行后处理，修复排版和常见识别错误"""
    if not text:
        return text
    
    import re
    
    # 【关键修复】先检测是否为代码 —— 如果是代码，直接委托给 _fix_code_ocr 做完整修复
    # 场景：上游 _is_code_text 可能漏判（如OCR原始输出格式变化），导致代码走入此函数
    # 此函数原有的段落合并逻辑对代码具有破坏性（会把多行压成一行）
    if _is_code_text(text):
        print("[DEBUG] _postprocess_ocr_text 检测到代码，委托给 _fix_code_ocr")
        return _fix_code_ocr(text)
    
    # 以下是非代码的普通文本路径（原有逻辑）
    # 0. 检测是否为英文为主内容（中文字符少于总字符的10%）
    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
    total_chars = len(text.replace(' ', '').replace('\n', ''))
    is_english_dominant = total_chars > 0 and chinese_chars / max(total_chars, 1) < 0.1
    
    # 1. 修复多余的空行：将连续3个以上换行压缩为2个
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 2. 核心排版修复：智能段落连接
    # 对于英文内容：如果一行以小写字母结尾且下一行以小写字母开头（或常见续行模式），
    # 说明是同一段落的断行，应连接；否则保持换行
    lines = text.split('\n')
    processed_lines = []
    i = 0
    while i < len(lines):
        line = lines[i]
        
        # 跳过空行（保留作为段落分隔符）
        if not line.strip():
            processed_lines.append('')
            i += 1
            continue
        
        # 尝试将当前行与后续非空行合并为段落
        paragraph_lines = [line]
        j = i + 1
        while j < len(lines):
            next_line = lines[j].strip()
            # 遇到空行，段落结束
            if not next_line:
                break
            prev_line = paragraph_lines[-1].strip()
            
            # 判断是否应该与下一行合并为同一段落：
            should_merge = False
            
            if is_english_dominant:
                # 英文规则：上一行以小写字母/逗号/分号结尾 且 下一行不以大写字母开头（除非是句子开头但明显在段内）
                prev_ends_lower = re.search(r'[a-z,;:\']$', prev_line)
                next_starts_lower_or_digit = re.match(r'^[a-z0-9\(\[\"\'—–]', next_line) or re.match(r'^[A-Z][a-z]*ing\b', next_line)
                
                # 或者下一行看起来像是前一句的延续（如 "and", "or", "but" 开头）
                next_is_continuation = re.match(r'^(and|or|but|so|for|nor|yet|which|that|who|whom|whose|when|where|while|as|if|than|into|onto|upon|from|with|within|without|about|above|below|between|among|through|during|before|after|since|until|against|along|across|behind|beyond|around|round|off|out|over|up|down|in|on|at|by|of|to|the|a|an|is|are|was|were|be|been|being|have|has|had|do|does|did|would|could|should|may|might|must|can|will|shall|it|its|this|these|those|their|they|he|she|we|you|I|me|him|her|us|them|my|your|his|our|their|not|no|also|just|only|even|still|already|very|too|more|most|much|such|same|other|another|all|any|both|each|every|few|many|several|some|what|which|who|how|why)\b', next_line, re.IGNORECASE)
                
                # 或者两行都较短（可能是同一行的断行）
                both_short = len(prev_line) < 80 and len(next_line) < 80
                
                should_merge = bool(prev_ends_lower and (next_starts_lower_or_digit or next_is_continuation)) or (
                    both_short and 
                    prev_ends_lower and
                    not re.match(r'^[A-Z][.!?]?\s*$', next_line) and
                    not next_line.endswith(('.', '!', '?', ':'))
                    # 排除标题/列表项特征
                    and not re.match(r'^\d+[\.\)]\s', next_line)
                    and not re.match(r'^[-•*]\s', next_line)
                    and not re.match(r'^#{1,6}\s', next_line)
                )
            else:
                # 中文/混合规则：一行末尾没有句号等结束标点，且下一行不是特殊格式，则可能需要连接
                prev_no_end_punct = not re.search(r'[。！？；：…」』】\)】》]', prev_line)
                next_not_special = not re.match(r'^[•\-\*\d#>\s【（《「『"』」》】】\]\)]', next_line)
                both_short = len(prev_line) < 60 and len(next_line) < 60
                
                should_merge = prev_no_end_punct and next_not_special and both_short
            
            if should_merge:
                paragraph_lines.append(next_line)
                j += 1
            else:
                break
        
        # 合并段落中的行为一段（用空格连接）
        merged_paragraph = ' '.join(l.strip() for l in paragraph_lines if l.strip())
        if merged_paragraph:
            processed_lines.append(merged_paragraph)
        
        # 跳过已处理的行
        if j > i + 1:
            # 添加中间遇到的空行
            for k in range(i + 1, j):
                if not lines[k].strip():
                    processed_lines.append('')
            i = j
        else:
            i += 1
    
    text = '\n'.join(processed_lines)
    
    # 3. 清理多余空格
    if is_english_dominant:
        # 英文文档：更激进地压缩空格
        text = re.sub(r' {2,}', ' ', text)
        # 标点前后多余空格
        text = re.sub(r'\s+([.,;:!?])', r'\1', text)
        # 确保句子后只有一个空格
        text = re.sub(r'([.!?])\s+', r'\1 ', text)
    
    # 4. 修复常见 OCR 错误字符映射
    text = re.sub(r'(?<=[\[\(\s,])@(?=[\]\)\s,])', '0', text)
    text = re.sub(r'(?<=[\d])@(?=[\d])', '0', text)
    text = re.sub(r'(?<=[\[\(\s,\d])®(?=[\]\)\s,\d])', '0', text)
    text = re.sub(r'y\s*=\s*\[([^\]]*)Q([^\]]*)\]', lambda m: 'y = [' + m.group(1) + '0' + m.group(2) + ']', text)
    text = re.sub(r'(X\s*=\s*)\n\s*\[', r'\1[', text)
    
    # 5. 最终去除首尾空白
    text = text.strip()
    
    return text


def _ocr_with_deepseek_vision(filepath, fid):
    """OCR 最终回退：尝试用 AI 视觉模型识别图片文字"""
    try:
        from PIL import Image

        # 读取图片并转为 base64
        img = Image.open(filepath)
        max_size = 2048
        if max(img.size) > max_size:
            ratio = max_size / max(img.size)
            img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.LANCZOS)

        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        img_b64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

        api_key, base_url = get_api_config()
        if api_key == 'sk-your-deepseek-api-key-here' or not api_key:
            return jsonify({"error": (
                "本地 OCR 引擎不可用。\n\n"
                "请安装 RapidOCR 以支持中文识别：\n"
                "   pip install rapidocr-onnxruntime\n\n"
                "然后重启应用即可。"
            )}), 400

        # 尝试调用视觉模型（需要 API 支持 vision/image 输入）
        # 注意: deepseek-chat 不支持图片！这里尝试调用，如果不支持会优雅降级
        import requests as req_lib
        resp = req_lib.post(
            f"{base_url.rstrip('/')}/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            },
            json={
                "model": "deepseek-chat",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{img_b64}"}
                            },
                            {
                                "type": "text",
                                "text": "请完整识别这张图片中的所有文字内容，保持原有的段落结构和格式。只输出识别出的文字，不要添加任何解释或描述。如果图片中没有可识别的文字，输出：[无文字内容]"
                            }
                        ]
                    }
                ],
                "temperature": 0.3,
                "max_tokens": 4096
            },
            timeout=60
        )

        if resp.status_code == 400:
            # API 不支持图片输入（如 deepseek-chat），直接给出安装指引
            return jsonify({"error": (
                "当前 API 不支持图片识别（vision）。\n\n"
                "要识别中文图片，请安装本地 OCR 引擎：\n"
                "   pip install rapidocr-onnxruntime\n\n"
                "安装后重启应用即可。"
            )}), 400

        resp.raise_for_status()
        data = resp.json()
        text = ""
        if "choices" in data and data["choices"]:
            text = data["choices"][0].get("message", {}).get("content", "")

        if text and len(text.strip()) > 2 and "[无文字内容]" not in text:
            db = get_db()
            db.execute("UPDATE files SET ocr_text = ? WHERE id = ?", (text.strip(), fid))
            db.commit()
            return jsonify({"text": text.strip(), "method": "ai-vision", "note": "AI视觉识别"})
        else:
            raise Exception("AI模型返回空结果")
    except Exception as e:
        err_msg = str(e)
        print(f"[OCR] Vision API 回退失败: {err_msg}")
        # 过滤掉不友好的错误信息
        if "400" in err_msg or "Bad Request" in err_msg or "image" in err_msg.lower():
            return jsonify({"error": (
                "本地 OCR 不可用，且当前 AI 模型不支持图片识别。\n\n"
                "解决方案（选一个）：\n"
                "1. 安装中文 OCR: pip install rapidocr-onnxruntime\n"
                "2. 安装 Tesseract-OCR 并勾选中文语言包\n"
                "3. 切换到支持视觉的 AI 模型（如 gpt-4o）\n\n"
                "安装后重启应用。"
            )}), 400
        return jsonify({
            "error": f"OCR 识别失败。\n\n{err_msg}\n\n"
                     "建议: pip install rapidocr-onnxruntime 并重启应用"
        }), 400

@app.route('/api/files/<fid>/summarize', methods=['POST'])
def api_file_summarize(fid):
    """AI 总结文档 - 分步调用确保每个部分完整输出"""
    db = get_db()
    row = db.execute("SELECT * FROM files WHERE id = ? AND is_deleted = 0 AND user_id = ?", (fid, get_current_user_id())).fetchone()
    if not row:
        return jsonify({"error": "文件不存在"}), 404
    
    if not row['ocr_text']:
        filepath = row['file_path']
        text = extract_text_from_file(filepath, row['file_type'])
        if text:
            db.execute("UPDATE files SET ocr_text = ? WHERE id = ?", (text, fid))
            db.commit()
    
    text = row['ocr_text'] or extract_text_from_file(row['file_path'], row['file_type'])
    if not text:
        return jsonify({"error": "无法提取文档内容"}), 400
    
    # 智能截断
    max_chars = 20000
    if len(text) > max_chars:
        head_len = int(max_chars * 0.7)
        tail_len = max_chars - head_len
        text = text[:head_len] + "\n\n...(中间内容已省略)...\n\n" + text[-tail_len:]
    
    prompt = f"""请对以下文档进行全面深入的分析总结，内容必须丰富详细，不能敷衍。

请自由组织你的总结，但务必包含以下方面的深入分析（不要写成固定章节标题，而是自然流畅地呈现）：

- 文档的核心议题、背景和目的（是什么、为什么）
- 关键的方法、过程或论证逻辑（怎么做）
- 核心发现、数据、结论（发现了什么，数据是什么）
- 文档的整体结构和各部分要点

具体要求：
- 每个方面都要写3-5句话以上，包含具体内容，不能一两句话带过
- 如果有具体数据（数值、百分比、时间等），必须明确列出
- 如果遇到数学公式，请按标准LaTeX原样输出变量符号（如 $x^2$、$\\frac{{a}}{{b}}$），不要解释变量为常数；不确定时保持原式，不要"脑补"
- 直接输出完整的总结报告，不要加"以下是对文档的分析"之类的前缀
- 用中文输出

文档内容：
{text}"""

    MATH_FIDELITY = ("你是一个专业的文档深度分析专家。"
                     "如果文档包含数学公式，请严格按照LaTeX语法输出（用 $...$ 包裹行内公式、$$...$$ 包裹块级公式），"
                     "保留原始数学符号，不要擅自简化或解释。")

    result = call_deepseek(
        [{"role": "system", "content": MATH_FIDELITY + "你的任务是全面深入地分析文档，输出丰富详细的总结报告。每个分析点都要有充分的具体内容，绝对不要写空泛的概括。"},
         {"role": "user", "content": prompt}],
        temperature=0.5, max_tokens=8192,
    )
    
    if "error" in result:
        return jsonify({"error": result["error"]}), 500
    
    summary = result["content"].strip()
    
    # 内容太短则重试
    if len(summary) < 300:
        retry = call_deepseek(
            [{"role": "system", "content": MATH_FIDELITY + "你必须输出非常详细完整的总结报告，字数不少于800字。如果输出简短将被视为失败。"},
             {"role": "user", "content": prompt + "\n\n【重要】上一次回复太简短，请务必输出详细、完整的分析，每个要点都要充分展开，总字数不少于800字。"}],
            temperature=0.7, max_tokens=8192,
        )
        if "error" not in retry and len(retry["content"].strip()) > len(summary):
            summary = retry["content"].strip()
    
    return jsonify({"summary": summary})

# ---------- 翻译 ----------

@app.route('/api/translate', methods=['POST'])
def api_translate():
    data = request.get_json()
    text = data.get('text', '')
    source_lang = data.get('source_lang', 'auto')
    target_lang = data.get('target_lang', 'zh')
    context = data.get('context', 'general')  # general/academic/business/casual
    
    if not text:
        return jsonify({"error": "请输入文本"}), 400
    
    context_prompts = {
        'academic': '请使用学术化的语言风格进行翻译，保持专业术语的准确性。',
        'business': '请使用商务正式的语言风格进行翻译。',
        'casual': '请使用自然口语化的语言风格进行翻译。',
        'general': '请准确流畅地翻译。',
        'screen_ocr': (
            '你是一个OCR截图翻译专家。请将以下OCR识别的文字翻译成目标语言。'
            '重要规则：\n'
            '1. 保持原文的段落结构、换行和排版格式，不要把所有内容堆在一起\n'
            '2. 保留原文中的编号、项目符号、缩进等格式\n'
            '3. 如果原文是代码，保留代码结构和关键字，只翻译注释和字符串\n'
            '4. 如果原文中有表格结构，尽量保持表格的行列对齐\n'
            '5. 只输出翻译后的文本，不要添加任何额外说明'
        )
    }
    
    system_prompt = f"你是一个专业的翻译助手。{context_prompts.get(context, context_prompts['general'])}"
    
    result = call_deepseek([
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"请将以下{source_lang}文本翻译成{target_lang}：\n\n{text}"}
    ])
    
    if "error" in result:
        return jsonify(result), 500
    
    translation = result['content']
    
    # 保存翻译记录
    db = get_db()
    tid = generate_id()
    db.execute(
        "INSERT INTO translations (id, source_text, translated_text, source_lang, target_lang, context_type, user_id) VALUES (?,?,?,?,?,?,?)",
        (tid, text, translation, source_lang, target_lang, context, get_current_user_id())
    )
    db.commit()
    
    return jsonify({
        "id": tid,
        "translation": translation,
        "source_text": text,
        "source_lang": source_lang,
        "target_lang": target_lang
    })

@app.route('/api/translate/batch', methods=['POST'])
def api_translate_batch():
    data = request.get_json()
    texts = data.get('texts', [])
    source_lang = data.get('source_lang', 'auto')
    target_lang = data.get('target_lang', 'zh')
    
    if not texts:
        return jsonify({"error": "请提供文本列表"}), 400
    
    combined = "\n---SEPARATOR---\n".join(texts)
    result = call_deepseek([
        {"role": "system", "content": "你是一个翻译助手。请翻译以下每段文本，用 ---SEPARATOR--- 分隔各段翻译结果。"},
        {"role": "user", "content": f"将以下{source_lang}文本翻译成{target_lang}，每段分别翻译，用 ---SEPARATOR--- 分隔：\n\n{combined}"}
    ])
    
    if "error" in result:
        return jsonify(result), 500
    
    translations = [t.strip() for t in result['content'].split('---SEPARATOR---')]
    return jsonify({"translations": translations, "source_texts": texts})

@app.route('/api/translations', methods=['GET'])
def api_get_translations():
    db = get_db()
    rows = db.execute(
        "SELECT * FROM translations WHERE user_id = ? ORDER BY created_at DESC LIMIT 200",
        (get_current_user_id(),)
    ).fetchall()
    results = []
    for r in rows:
        d = dict(r)
        d['tags'] = json.loads(d['tags']) if d['tags'] else []
        results.append(d)
    return jsonify(results)

@app.route('/api/translations/<tid>', methods=['DELETE'])
def api_delete_translation(tid):
    db = get_db()
    db.execute("DELETE FROM translations WHERE id = ?", (tid,))
    db.commit()
    return jsonify({"message": "已删除"})

@app.route('/api/translations/export', methods=['POST'])
def api_export_translations():
    data = request.get_json()
    ids = data.get('ids', [])
    fmt = data.get('format', 'txt')
    
    db = get_db()
    translations = []
    for tid in ids:
        row = db.execute("SELECT * FROM translations WHERE id = ? AND user_id = ?", (tid, get_current_user_id())).fetchone()
        if row:
            translations.append(dict(row))
    
    if not translations:
        return jsonify({"error": "没有找到翻译记录"}), 404
    
    if fmt == 'txt':
        content = ""
        for t in translations:
            content += f"原文 ({t['source_lang']}):\n{t['source_text']}\n\n"
            content += f"译文 ({t['target_lang']}):\n{t['translated_text']}\n"
            content += f"时间: {t['created_at']}\n"
            content += "=" * 50 + "\n\n"
        
        filepath = Path(app.config['UPLOAD_FOLDER']) / f"translation_export_{generate_id()}.txt"
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        return send_file(filepath, as_attachment=True, download_name="翻译导出.txt")
    
    elif fmt == 'json':
        filepath = Path(app.config['UPLOAD_FOLDER']) / f"translation_export_{generate_id()}.json"
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(translations, f, ensure_ascii=False, indent=2)
        return send_file(filepath, as_attachment=True, download_name="翻译导出.json")
    
    return jsonify({"error": "不支持的格式"}), 400

# ---------- AI 聊天 / 智能体 ----------

@app.route('/api/ai/chats', methods=['GET'])
def api_get_chats():
    db = get_db()
    rows = db.execute("SELECT id, title, created_at, updated_at FROM ai_chats WHERE user_id = ? ORDER BY updated_at DESC", (get_current_user_id(),)).fetchall()
    return jsonify([dict(r) for r in rows])

@app.route('/api/ai/chats', methods=['POST'])
def api_create_chat():
    data = request.get_json()
    cid = generate_id()
    db = get_db()
    db.execute(
        "INSERT INTO ai_chats (id, title, messages, user_id) VALUES (?,?,?,?)",
        (cid, data.get('title', '新对话'), json.dumps([], ensure_ascii=False), get_current_user_id())
    )
    db.commit()
    return jsonify({"id": cid})

@app.route('/api/ai/chats/<cid>', methods=['GET'])
def api_get_chat(cid):
    db = get_db()
    row = db.execute("SELECT * FROM ai_chats WHERE id = ? AND user_id = ?", (cid, get_current_user_id())).fetchone()
    if not row:
        return jsonify({"error": "对话不存在"}), 404
    chat = dict(row)
    chat['messages'] = json.loads(chat['messages']) if chat['messages'] else []
    return jsonify(chat)

@app.route('/api/ai/chats/<cid>', methods=['DELETE'])
def api_delete_chat(cid):
    db = get_db()
    db.execute("DELETE FROM ai_chats WHERE id = ?", (cid,))
    db.commit()
    return jsonify({"message": "已删除"})

@app.route('/api/ai/chat', methods=['POST'])
def api_ai_chat():
    """AI 对话接口"""
    data = request.get_json()
    chat_id = data.get('chat_id')
    message = data.get('message', '')
    mode = data.get('mode', 'chat')  # chat/translate/summarize/quiz/analyze
    
    if not message:
        return jsonify({"error": "请输入消息"}), 400
    
    # 数学公式保真约束（追加到所有模式）
    MATH_FIDELITY_TAIL = (" 如果回复涉及数学公式，请严格按LaTeX语法输出："
                          "行内公式用 $...$ 包裹，块级公式用 $$...$$ 包裹；"
                          "保留原始变量符号（如 $\\sigma$、$\\partial$、$\\int$），"
                          "不要将希腊字母/运算符解释为常数，不确定时保持原式不要编造。")
    
    # 构建系统提示词
    system_prompts = {
        'chat': "你是一个全能的AI助手，可以帮助用户解答各类问题、分析文档、提供建议。请用中文回复。",
        'translate': "你是一个专业的翻译专家，精通中英互译及其他多语种翻译。请提供准确流畅的翻译。",
        'summarize': "你是一个文档分析专家，擅长总结长文档的核心内容、提炼要点、生成大纲和思维导图。",
        'quiz': "你是一个教育评测专家，擅长根据学习内容生成高质量的练习题。",
        'analyze': "你是一个数据分析专家，擅长从文档中提取关键信息、数据和分析结论。"
    }
    
    system_msg = system_prompts.get(mode, system_prompts['chat']) + MATH_FIDELITY_TAIL
    
    # 获取对话历史
    db = get_db()
    conversation = [{"role": "system", "content": system_msg}]
    
    if chat_id:
        row = db.execute("SELECT messages FROM ai_chats WHERE id = ? AND user_id = ?", (chat_id, get_current_user_id())).fetchone()
        if row:
            history = json.loads(row['messages']) if row['messages'] else []
            # 取最近20轮对话
            conversation.extend(history[-40:])
    
    conversation.append({"role": "user", "content": message})
    
    # 根据模式增强 prompt
    if mode == 'quiz':
        conversation[-1]["content"] = f"请根据以下内容生成3-5道练习题（包含选择题、填空题、简答题），并附带答案和解析：\n\n{message}"
    elif mode == 'summarize':
        conversation[-1]["content"] = f"请总结以下内容的核心要点、生成结构化大纲，并提炼关键信息：\n\n{message}"
    
    result = call_deepseek(conversation, temperature=0.7, max_tokens=8192)
    
    if "error" in result:
        return jsonify(result), 500
    
    # 更新或创建对话
    if not chat_id:
        chat_id = generate_id()
        title = message[:30] + ('...' if len(message) > 30 else '')
        db.execute(
            "INSERT INTO ai_chats (id, title, messages, user_id) VALUES (?,?,?,?)",
            (chat_id, title, json.dumps([], ensure_ascii=False), get_current_user_id())
        )
    
    # 更新消息历史
    row = db.execute("SELECT messages FROM ai_chats WHERE id = ? AND user_id = ?", (chat_id, get_current_user_id())).fetchone()
    if row:
        messages = json.loads(row['messages']) if row['messages'] else []
        messages.append({"role": "user", "content": message})
        messages.append({"role": "assistant", "content": result['content']})
        db.execute(
            "UPDATE ai_chats SET messages = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
            (json.dumps(messages, ensure_ascii=False), chat_id)
        )
        db.commit()
    
    return jsonify({
        "chat_id": chat_id,
        "reply": result['content'],
        "mode": mode
    })

# ---------- 配置 ----------

@app.route('/api/config', methods=['GET'])
def api_get_config():
    db = get_db()
    rows = db.execute("SELECT key, value FROM config").fetchall()
    config = {r['key']: r['value'] for r in rows}
    # 隐藏 API key 的部分内容
    if config.get('deepseek_api_key'):
        key = config['deepseek_api_key']
        if len(key) > 8:
            config['deepseek_api_key_masked'] = key[:4] + '*' * (len(key) - 8) + key[-4:]
    return jsonify(config)

@app.route('/api/config', methods=['POST'])
def api_update_config():
    data = request.get_json()
    db = get_db()
    for key, value in data.items():
        db.execute(
            "INSERT INTO config (key, value) VALUES (?,?) ON CONFLICT(key) DO UPDATE SET value = ?",
            (key, value, value)
        )
    db.commit()
    return jsonify({"message": "配置已更新"})

# ---------- 批量处理 ----------

@app.route('/api/batch/process', methods=['POST'])
def api_batch_process():
    """批量处理文件"""
    data = request.get_json()
    action = data.get('action', '')
    file_ids = data.get('file_ids', [])
    
    if not file_ids:
        return jsonify({"error": "请选择文件"}), 400
    
    db = get_db()
    results = []
    
    for fid in file_ids:
        row = db.execute("SELECT * FROM files WHERE id = ? AND is_deleted = 0", (fid,)).fetchone()
        if not row:
            results.append({"id": fid, "error": "文件不存在"})
            continue
        
        filepath = row['file_path']
        if not os.path.exists(filepath):
            results.append({"id": fid, "error": "文件已丢失"})
            continue
        
        if action == 'extract_text':
            text = extract_text_from_file(filepath, row['file_type'])
            if text:
                db.execute("UPDATE files SET ocr_text = ? WHERE id = ?", (text, fid))
                db.commit()
            results.append({"id": fid, "name": row['original_name'], "text_preview": text[:200] if text else ""})
        
        elif action == 'summarize':
            text = row['ocr_text'] or extract_text_from_file(filepath, row['file_type'])
            if text:
                text = text[:10000]
                result = call_deepseek([
                    {"role": "system", "content": "请简洁总结以下文档的核心内容（200字以内）。"},
                    {"role": "user", "content": text}
                ])
                summary = result.get('content', result.get('error', '总结失败'))
                results.append({"id": fid, "name": row['original_name'], "summary": summary})
            else:
                results.append({"id": fid, "name": row['original_name'], "error": "无法提取文本"})
    
    return jsonify({"results": results})

@app.route('/api/batch/compare', methods=['POST'])
def api_batch_compare():
    """跨文件对比分析"""
    data = request.get_json()
    file_ids = data.get('file_ids', [])
    
    if len(file_ids) < 2:
        return jsonify({"error": "对比分析需要至少2个文件"}), 400
    
    db = get_db()
    
    # 收集所有文件文本
    file_texts = []
    for fid in file_ids:
        row = db.execute("SELECT * FROM files WHERE id = ? AND is_deleted = 0", (fid,)).fetchone()
        if not row:
            continue
        text = row['ocr_text'] or extract_text_from_file(row['file_path'], row['file_type'])
        # 截取每份文档最多8000字符，总文档不超过25000
        truncated_text = text[:8000] if text else ''
        file_texts.append({
            'name': row['original_name'],
            'text': truncated_text,
            'type': row['file_type']
        })
    
    if not file_texts:
        return jsonify({"error": "无法提取任何文件的文本内容"}), 400
    
    if len(file_texts) < 2:
        return jsonify({"error": "有效文件不足2个"}), 400
    
    # 拼接所有文档内容
    combined_docs = '\n\n========== 文档分隔线 ==========\n\n'.join(
        f"【文档 {i+1}: {f['name']}】\n{f['text']}" for i, f in enumerate(file_texts)
    )
    
    # 如果总文本太长，截断
    max_total = 25000
    if len(combined_docs) > max_total:
        combined_docs = combined_docs[:max_total] + "\n\n（部分内容已省略）"
    
    prompt = f"""你是一个专业的多文档对比分析专家。请对以下 {len(file_texts)} 份文档进行深入的对比分析。

## 要求
请从以下维度进行全面分析，每个维度都要有具体的、基于原文内容的详细论述：

### 1. 各文档核心主题概述
逐个简要概括每份文档的核心主题、主要内容和关键信息。

### 2. 关联性分析
这些文档之间有什么关联？是否存在：
- 相互引用或讨论同一主题？
- 时间上的先后或逻辑关系？
- 共同涉及的人物、事件、数据、理论？

### 3. 差异性分析
这些文档之间的主要差异是什么？包括但不限于：
- 观点差异（是否持不同立场？）
- 方法论差异（研究方法/处理方式不同？）
- 结论差异（得出的结论是否一致或有冲突？）
- 数据来源和可信度差异

### 4. 综合评价与洞察
基于以上分析，给出你的综合评价。如果把这些文档放在一起看，能得出什么新的结论或发现？

## 重要要求
- 必须引用具体内容来支撑你的分析（如："文档A指出X，而文档B认为Y"）
- 不要泛泛而谈，要深入到具体细节
- 如有数学公式，请严格按LaTeX输出（$...$ 行内、$$...$$ 块级），不要编造或简化
- 用中文输出

{combined_docs}"""

    COMPARE_MATH = ("你是专业的多文档对比分析专家。"
                    "如果文档包含数学公式，请严格按照LaTeX语法输出，保留原始数学符号。")

    result = call_deepseek(
        [{"role": "system", "content": COMPARE_MATH + "你需要对多个文档进行深入比较，找出它们之间的关联性、差异性，并给出综合洞察。每个分析点都必须有具体内容支撑。"},
         {"role": "user", "content": prompt}],
        temperature=0.3, max_tokens=16384
    )
    
    if "error" in result:
        return jsonify({"error": result["error"]}), 500
    
    comparison = result['content'].strip()
    
    # 验证结果长度 - 如果太短说明可能被截断，用更高temperature重试一次
    if len(comparison) < 500:
        result2 = call_deepseek(
            [{"role": "system", "content": COMPARE_MATH + "你必须输出非常详细的对比分析报告，不能省略任何部分。"},
             {"role": "user", "content": prompt + "\n\n【重要提示】请务必完整输出所有4个维度的分析，每个维度都要详细论述，不要省略或概括。"}],
            temperature=0.7, max_tokens=16384
        )
        if "error" not in result2 and len(result2["content"].strip()) > len(comparison):
            comparison = result2["content"].strip()
    
    return jsonify({"comparison": comparison})

# ---------- 搜索 ----------

@app.route('/api/search', methods=['GET'])
def api_fulltext_search():
    q = request.args.get('q', '')
    if not q or len(q) < 2:
        return jsonify({"notes": [], "files": [], "translations": []})
    
    db = get_db()
    uid = get_current_user_id()
    notes = db.execute(
        "SELECT id, title, plain_text FROM notes WHERE is_deleted=0 AND user_id = ? AND (title LIKE ? OR plain_text LIKE ?) LIMIT 20",
        (uid, f'%{q}%', f'%{q}%')
    ).fetchall()
    files = db.execute(
        "SELECT id, original_name, ocr_text FROM files WHERE is_deleted=0 AND user_id = ? AND (original_name LIKE ? OR ocr_text LIKE ?) LIMIT 20",
        (uid, f'%{q}%', f'%{q}%')
    ).fetchall()
    translations = db.execute(
        "SELECT id, source_text, translated_text FROM translations WHERE user_id = ? AND (source_text LIKE ? OR translated_text LIKE ?) LIMIT 20",
        (uid, f'%{q}%', f'%{q}%')
    ).fetchall()
    return jsonify({
        "notes": [dict(r) for r in notes],
        "files": [dict(r) for r in files],
        "translations": [dict(r) for r in translations]
    })

# ---------- 主页 ----------

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/stats', methods=['GET'])
def api_stats():
    db = get_db()
    uid = get_current_user_id()
    note_count = db.execute("SELECT COUNT(*) FROM notes WHERE is_deleted=0 AND user_id = ?", (uid,)).fetchone()[0]
    file_count = db.execute("SELECT COUNT(*) FROM files WHERE is_deleted=0 AND user_id = ?", (uid,)).fetchone()[0]
    translation_count = db.execute("SELECT COUNT(*) FROM translations WHERE user_id = ?", (uid,)).fetchone()[0]
    chat_count = db.execute("SELECT COUNT(*) FROM ai_chats WHERE user_id = ?", (uid,)).fetchone()[0]
    
    return jsonify({
        "notes": note_count,
        "files": file_count,
        "translations": translation_count,
        "chats": chat_count,
    })

# ==================== 用户登录/注册 ====================

import smtplib
import random
from email.mime.text import MIMEText

def login_required(f):
    """登录验证装饰器"""
    @wraps(f)
    def decorated(*args, **kwargs):
        if 'user_id' not in session:
            return jsonify({"error": "请先登录"}), 401
        return f(*args, **kwargs)
    return decorated

def get_current_user():
    """获取当前登录用户"""
    if 'user_id' not in session:
        return None
    db = get_db()
    return db.execute("SELECT * FROM users WHERE id = ?", (session['user_id'],)).fetchone()

def get_current_user_id():
    """获取当前用户ID，未登录返回 -1（不存在的ID，查询结果为空）"""
    return session.get('user_id', -1)

@app.route('/api/auth/send-code', methods=['POST'])
def api_send_code():
    """发送邮箱验证码"""
    data = request.get_json()
    email = (data.get('email') or '').strip()
    if not email or '@' not in email:
        return jsonify({"error": "请输入有效的邮箱地址"}), 400

    code = ''.join(random.choices('0123456789', k=6))
    db = get_db()
    db.execute("DELETE FROM captcha_codes WHERE email = ?", (email,))
    db.execute("INSERT INTO captcha_codes (email, code) VALUES (?, ?)", (email, code))
    db.commit()

    # 发送邮件
    try:
        db2 = get_db()
        configs = {r['key']: r['value'] for r in db2.execute("SELECT key, value FROM config").fetchall()}
        smtp_host = configs.get('smtp_host', 'smtp.qq.com')
        smtp_port = int(configs.get('smtp_port', '587'))
        smtp_user = configs.get('smtp_user', '')
        smtp_pass = configs.get('smtp_pass', '')
        if smtp_user and smtp_pass:
            msg = MIMEText(f'您的验证码是：{code}\n\n验证码5分钟内有效，请勿泄露给他人。', 'plain', 'utf-8')
            msg['Subject'] = 'DocMind 登录验证码'
            msg['From'] = smtp_user
            msg['To'] = email
            with smtplib.SMTP(smtp_host, smtp_port) as server:
                server.starttls()
                server.login(smtp_user, smtp_pass)
                server.send_message(msg)
            return jsonify({"message": "验证码已发送"})
    except Exception as e:
        print(f"[AUTH] 邮件发送失败: {e}")

    print(f"[AUTH] 调试验证码 [{code}] -> {email}")
    return jsonify({"message": "（调试模式）验证码已生成", "debug_code": code,
                    "hint": "请在设置中配置 SMTP 邮箱以正常发送邮件"})

@app.route('/api/auth/register', methods=['POST'])
def api_register():
    """注册（用户名+密码，无需邮箱验证码）"""
    data = request.get_json()
    username = (data.get('username') or '').strip()
    password = data.get('password', '')
    if not username:
        return jsonify({"error": "请输入用户名"}), 400
    if len(password) < 6:
        return jsonify({"error": "密码至少6位"}), 400
    db = get_db()
    if db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone():
        return jsonify({"error": "该用户名已被注册"}), 400
    password_hash = generate_password_hash(password)
    db.execute("INSERT INTO users (username, password_hash, verified) VALUES (?, ?, 1)",
               (username, password_hash))
    db.commit()
    user = db.execute("SELECT * FROM users WHERE username = ?", (username,)).fetchone()
    session['user_id'] = user['id']
    session['username'] = user['username']
    session.permanent = True
    return jsonify({"message": "注册成功", "user": {"id": user['id'], "username": user['username'], "email": ''}})

@app.route('/api/auth/login', methods=['POST'])
def api_login():
    """登录（用户名+密码）"""
    data = request.get_json()
    username = (data.get('username') or '').strip()
    password = data.get('password', '')
    db = get_db()
    # 支持用户名或邮箱登录
    user = db.execute("SELECT * FROM users WHERE username = ? OR email = ?", (username, username)).fetchone()
    if not user or not check_password_hash(user['password_hash'], password):
        return jsonify({"error": "用户名或密码错误"}), 401
    session['user_id'] = user['id']
    session['username'] = user['username']
    session.permanent = True
    return jsonify({"message": "登录成功", "user": {"id": user['id'], "username": user['username']}})

@app.route('/api/auth/logout', methods=['POST'])
def api_logout():
    session.clear()
    return jsonify({"message": "已退出登录"})

@app.route('/api/auth/me', methods=['GET'])
def api_auth_me():
    if 'user_id' not in session:
        return jsonify({"user": None})
    user = get_current_user()
    if not user:
        session.clear()
        return jsonify({"user": None})
    return jsonify({"user": {"id": user['id'], "username": user['username']}})

if __name__ == '__main__':
    app.run(debug=False, host='0.0.0.0', port=5000)
